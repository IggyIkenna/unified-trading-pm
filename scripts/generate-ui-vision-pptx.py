#!/usr/bin/env python3
"""Generate the Unified Trading Platform UI/UX Redesign Vision PowerPoint.

Uses real strategy data, SMA model, and actual P&L components from the system.
Run: python3 unified-trading-pm/scripts/generate-ui-vision-pptx.py
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── Design tokens ──────────────────────────────────────────────────────────────

BG_PRIMARY = RGBColor(0x0A, 0x0A, 0x0B)
BG_SECONDARY = RGBColor(0x11, 0x11, 0x13)
BG_TERTIARY = RGBColor(0x18, 0x18, 0x1B)
BG_ELEVATED = RGBColor(0x1C, 0x1C, 0x1F)
BG_CARD = RGBColor(0x1E, 0x1E, 0x22)

TEXT_PRIMARY = RGBColor(0xFA, 0xFA, 0xFA)
TEXT_SECONDARY = RGBColor(0xA1, 0xA1, 0xAA)
TEXT_TERTIARY = RGBColor(0x71, 0x71, 0x7A)
TEXT_MUTED = RGBColor(0x52, 0x52, 0x5B)

ACCENT_CYAN = RGBColor(0x22, 0xD3, 0xEE)
ACCENT_GREEN = RGBColor(0x4A, 0xDE, 0x80)
ACCENT_AMBER = RGBColor(0xFB, 0xBF, 0x24)
ACCENT_RED = RGBColor(0xF8, 0x71, 0x71)
ACCENT_BLUE = RGBColor(0x60, 0xA5, 0xFA)
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
ACCENT_TEAL = RGBColor(0x2D, 0xD4, 0xBF)
ACCENT_ORANGE = RGBColor(0xF9, 0x73, 0x16)

BORDER = RGBColor(0x27, 0x27, 0x2A)
BORDER_SUBTLE = RGBColor(0x1F, 0x1F, 0x23)

FONT_SANS = "IBM Plex Sans"
FONT_MONO = "JetBrains Mono"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color=BG_PRIMARY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=12,
    font_name=FONT_SANS,
    color=TEXT_PRIMARY,
    bold=False,
    alignment=PP_ALIGN.LEFT,
    font_color_override=None,
):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.color.rgb = font_color_override or color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_card(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Subtle border
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.5)
    shape.rotation = 0.0
    return shape


def add_pill(slide, left, top, width, text, color=ACCENT_CYAN, font_size=9):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.28))
    shape.fill.solid()
    # Dim background version of the color
    min(color[0] + 10, 255)
    min(color[1] + 10, 255)
    min(color[2] + 10, 255)
    shape.fill.fore_color.rgb = RGBColor(max(color[0] // 5, 15), max(color[1] // 5, 15), max(color[2] // 5, 15))
    shape.line.color.rgb = RGBColor(
        min(color[0] // 2 + 30, 200), min(color[1] // 2 + 30, 200), min(color[2] // 2 + 30, 200)
    )
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = FONT_SANS
    p.font.color.rgb = color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_sparkline(slide, left, top, width, height, color, pattern="up"):
    """Draw a simple sparkline shape using connected line segments."""
    if pattern == "up":
        points = [(0, 0.7), (0.15, 0.6), (0.3, 0.65), (0.45, 0.4), (0.6, 0.5), (0.75, 0.3), (0.9, 0.15), (1.0, 0.1)]
    elif pattern == "down":
        points = [(0, 0.2), (0.15, 0.3), (0.3, 0.25), (0.45, 0.5), (0.6, 0.4), (0.75, 0.6), (0.9, 0.7), (1.0, 0.8)]
    elif pattern == "flat":
        points = [(0, 0.5), (0.15, 0.45), (0.3, 0.55), (0.45, 0.5), (0.6, 0.48), (0.75, 0.52), (0.9, 0.5), (1.0, 0.49)]
    else:  # volatile
        points = [(0, 0.5), (0.15, 0.3), (0.3, 0.7), (0.45, 0.2), (0.6, 0.6), (0.75, 0.35), (0.9, 0.55), (1.0, 0.4)]
    # Draw as small dots (shapes can't do polylines easily, use tiny rectangles)
    for px, py in points:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(left + px * width), Inches(top + py * height), Inches(0.04), Inches(0.04)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()


def add_kpi_card(slide, left, top, label, value, sub, accent=ACCENT_CYAN, sparkline_pattern="up"):
    add_card(slide, left, top, 2.0, 1.3, BG_ELEVATED)
    add_text_box(slide, left + 0.15, top + 0.1, 1.7, 0.25, label, font_size=9, color=TEXT_TERTIARY)
    add_text_box(
        slide,
        left + 0.15,
        top + 0.35,
        1.7,
        0.35,
        value,
        font_size=22,
        font_name=FONT_MONO,
        bold=True,
        font_color_override=accent,
    )
    add_text_box(slide, left + 0.15, top + 0.78, 1.7, 0.25, sub, font_size=8, color=TEXT_MUTED, font_name=FONT_MONO)
    # Sparkline under the value
    add_sparkline(slide, left + 0.15, top + 0.95, 1.6, 0.25, accent, sparkline_pattern)


def add_nav_bar(slide, active_idx=0, show_as_of=False, as_of_text="Live"):
    """Add the global nav bar at the top."""
    add_card(slide, 0, 0, 13.333, 0.42, BG_SECONDARY)
    labels = [
        ("Trading", ACCENT_GREEN),
        ("Strategy", ACCENT_BLUE),
        ("Markets", ACCENT_PURPLE),
        ("Ops", ACCENT_AMBER),
        ("Config", ACCENT_CYAN),
        ("ML", ACCENT_ORANGE),
        ("Reports", ACCENT_TEAL),
    ]
    x = 0.3
    for i, (label, color) in enumerate(labels):
        c = color if i == active_idx else TEXT_MUTED
        add_text_box(slide, x, 0.08, 1.0, 0.3, label, font_size=10, bold=(i == active_idx), font_color_override=c)
        x += 1.15
    if show_as_of:
        add_pill(
            slide, 8.8, 0.08, 1.5, f"As-of: {as_of_text}", ACCENT_AMBER if as_of_text != "Live" else ACCENT_GREEN, 8
        )
    add_text_box(slide, 10.5, 0.08, 2.5, 0.3, "Search entity, run, service...", font_size=9, color=TEXT_MUTED)


def add_lifecycle_rail(slide, active_idx=3, top=0.48):
    """Add the lifecycle rail below nav."""
    steps = ["Design", "Simulate", "Promote", "Run", "Monitor", "Explain", "Reconcile"]
    x = 0.3
    for i, step in enumerate(steps):
        if i == active_idx:
            add_card(slide, x, top, 1.55, 0.35, RGBColor(0xFA, 0xFA, 0xFA))
            add_text_box(
                slide, x + 0.05, top + 0.02, 0.3, 0.15, f"0{i + 1}", font_size=7, font_color_override=BG_PRIMARY
            )
            add_text_box(
                slide, x + 0.05, top + 0.12, 1.4, 0.2, step, font_size=11, bold=True, font_color_override=BG_PRIMARY
            )
        else:
            add_card(slide, x, top, 1.55, 0.35, BG_ELEVATED)
            add_text_box(slide, x + 0.05, top + 0.02, 0.3, 0.15, f"0{i + 1}", font_size=7, color=TEXT_MUTED)
            add_text_box(slide, x + 0.05, top + 0.12, 1.4, 0.2, step, font_size=11, color=TEXT_SECONDARY)
        x += 1.72


def add_breadcrumb(slide, parts, top=0.9):
    text = "  >  ".join(parts)
    add_text_box(slide, 0.3, top, 10, 0.25, text, font_size=9, font_name=FONT_MONO, color=TEXT_TERTIARY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════


def slide_01_hero(prs):
    """Title + Command Center hero with real data."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Nav + lifecycle + breadcrumb
    add_nav_bar(slide, active_idx=0, show_as_of=True, as_of_text="Live")
    add_lifecycle_rail(slide, active_idx=3, top=0.48)
    add_breadcrumb(slide, ["Odum Delta One", "Apex Capital", "DEFI_ETH_STAKED_BASIS_SCE_1H", "cfg-2.4.0"], top=0.9)

    # Title block
    add_text_box(slide, 0.3, 1.2, 8, 0.3, "Unified Trading Platform", font_size=28, bold=True, color=TEXT_PRIMARY)
    add_text_box(slide, 0.3, 1.6, 8, 0.3, "UI / UX Redesign Vision", font_size=14, color=TEXT_SECONDARY)
    add_text_box(
        slide,
        0.3,
        1.95,
        8,
        0.25,
        "From 11 disconnected tools to 4 daily-use surfaces + 3 specialist tools.",
        font_size=10,
        color=TEXT_TERTIARY,
    )

    # Badges
    add_pill(slide, 9.0, 1.25, 1.6, "Plan 2026-03-17", ACCENT_CYAN, 8)
    add_pill(slide, 10.7, 1.25, 0.7, "Draft", ACCENT_AMBER, 8)
    add_pill(slide, 11.5, 1.25, 1.5, "Architecture", ACCENT_BLUE, 8)

    # KPI cards with sparklines
    add_kpi_card(slide, 0.3, 2.4, "Firm PnL", "+$1.42m", "  +0.8% 1d", ACCENT_GREEN, "up")
    add_kpi_card(slide, 2.5, 2.4, "Net Exposure", "$4.2m", "1.2x levered", ACCENT_CYAN, "volatile")
    add_kpi_card(slide, 4.7, 2.4, "Margin Health", "82%", "$340k free", ACCENT_AMBER, "flat")
    add_kpi_card(slide, 6.9, 2.4, "Live Strategies", "12", "9 ok, 2 warn, 1 paused", ACCENT_BLUE, "flat")
    add_kpi_card(slide, 9.1, 2.4, "Alerts", "3 crit", "2 high", ACCENT_RED, "down")

    # Strategy performance table
    add_card(slide, 0.3, 3.9, 6.5, 3.1, BG_ELEVATED)
    add_text_box(slide, 0.5, 4.0, 4, 0.25, "Strategy Performance", font_size=12, bold=True, color=TEXT_PRIMARY)
    add_text_box(
        slide,
        0.5,
        4.25,
        6,
        0.2,
        "Click strategy name -> analytics  |  Click status -> filtered positions",
        font_size=8,
        color=TEXT_MUTED,
    )

    # Table headers
    cols = [
        ("Strategy", 0.5, 2.2),
        ("Arch", 2.7, 0.6),
        ("St", 3.3, 0.4),
        ("PnL", 3.8, 0.8),
        ("Sharpe", 4.6, 0.6),
        ("DD", 5.3, 0.5),
        ("Trend", 5.9, 0.5),
    ]
    for label, x, w in cols:
        add_text_box(slide, x, 4.5, w, 0.2, label, font_size=8, bold=True, color=TEXT_MUTED)

    # Table rows — REAL strategies
    rows = [
        ("DEFI_ETH_STAKED_BASIS", "Yield", ACCENT_GREEN, "L", "+$289k", "2.5", "3.3%"),
        ("DEFI_ETH_BASIS_SCE_1H", "Arb", ACCENT_GREEN, "L", "+$412k", "2.1", "4.1%"),
        ("DEFI_USDT_LENDING_SCE", "Yield", ACCENT_GREEN, "L", "+$91k", "1.8", "2.1%"),
        ("DEFI_ETH_RECURSIVE", "Yield", ACCENT_AMBER, "W", "+$188k", "1.9", "5.2%"),
        ("TRADFI_SPY_ML_DIR_1H", "Dir", ACCENT_GREEN, "L", "+$67k", "1.4", "3.9%"),
        ("SPORTS_HT_ML_V2.1", "Dir", ACCENT_GREEN, "L", "+$44k", "1.6", "1.8%"),
        ("SPORTS_ARB", "Arb", ACCENT_GREEN, "L", "+$31k", "1.3", "2.4%"),
    ]
    y = 4.75
    for name, arch, status_color, st, pnl, sharpe, dd in rows:
        add_text_box(slide, 0.5, y, 2.2, 0.2, name, font_size=9, font_name=FONT_MONO, color=ACCENT_CYAN)
        add_text_box(slide, 2.7, y, 0.6, 0.2, arch, font_size=9, color=TEXT_TERTIARY)
        add_text_box(slide, 3.3, y, 0.4, 0.2, f"  {st}", font_size=9, font_color_override=status_color)
        pnl_color = ACCENT_GREEN if pnl.startswith("+") else ACCENT_RED
        add_text_box(slide, 3.8, y, 0.8, 0.2, pnl, font_size=9, font_name=FONT_MONO, font_color_override=pnl_color)
        add_text_box(slide, 4.6, y, 0.6, 0.2, sharpe, font_size=9, font_name=FONT_MONO, color=TEXT_SECONDARY)
        add_text_box(slide, 5.3, y, 0.5, 0.2, dd, font_size=9, font_name=FONT_MONO, color=TEXT_SECONDARY)
        y += 0.35

    # P&L + Risk Attribution panel
    add_card(slide, 7.0, 3.9, 3.0, 3.1, BG_ELEVATED)
    add_text_box(slide, 7.2, 4.0, 2.8, 0.25, "P&L + Risk Attribution", font_size=11, bold=True, color=TEXT_PRIMARY)
    add_text_box(slide, 7.2, 4.25, 2.8, 0.2, "Same 6D, two sides", font_size=8, color=TEXT_MUTED)

    attrib_rows = [
        ("Funding", "+$412k", "$8.2m"),
        ("Basis", "+$355k", "14 bps"),
        ("Staking Yield", "+$145k", "LTV .72"),
        ("Delta", "+$61k", "$2.4m net"),
        ("Greeks", "-$8k", "Delta:-0.98"),
        ("Slippage", "-$61k", "---"),
        ("Gas/Fees", "-$44k", "---"),
        ("Recon Pending", "-$18k", "4 breaks"),
    ]
    add_text_box(slide, 7.2, 4.5, 1.2, 0.2, "Component", font_size=8, bold=True, color=TEXT_MUTED)
    add_text_box(slide, 8.4, 4.5, 0.7, 0.2, "P&L", font_size=8, bold=True, color=TEXT_MUTED)
    add_text_box(slide, 9.1, 4.5, 0.8, 0.2, "Exposure", font_size=8, bold=True, color=TEXT_MUTED)
    y = 4.75
    for comp, pnl, exp in attrib_rows:
        add_text_box(slide, 7.2, y, 1.2, 0.18, comp, font_size=8, color=TEXT_SECONDARY)
        pnl_c = ACCENT_GREEN if pnl.startswith("+") else ACCENT_RED
        add_text_box(slide, 8.4, y, 0.7, 0.18, pnl, font_size=8, font_name=FONT_MONO, font_color_override=pnl_c)
        add_text_box(slide, 9.1, y, 0.8, 0.18, exp, font_size=8, font_name=FONT_MONO, color=TEXT_TERTIARY)
        y += 0.3

    # Alerts panel
    add_card(slide, 10.2, 3.9, 2.85, 1.5, BG_ELEVATED)
    add_text_box(slide, 10.35, 4.0, 2.5, 0.25, "Alerts & Incidents", font_size=11, bold=True, color=TEXT_PRIMARY)
    alerts = [
        ("CRIT", "Kill switch: DEFI_ETH_BASIS", ACCENT_RED),
        ("HIGH", "Feature freshness 92s lag EU", ACCENT_AMBER),
        ("MED", "Recon break: Apex SMA", ACCENT_AMBER),
    ]
    y = 4.35
    for sev, msg, color in alerts:
        add_text_box(slide, 10.35, y, 0.5, 0.2, sev, font_size=8, bold=True, font_color_override=color)
        add_text_box(slide, 10.85, y, 2.1, 0.2, msg, font_size=8, color=TEXT_SECONDARY)
        y += 0.35

    # Health panel
    add_card(slide, 10.2, 5.6, 2.85, 1.4, BG_ELEVATED)
    add_text_box(slide, 10.35, 5.7, 2.5, 0.25, "Health & Freshness", font_size=11, bold=True, color=TEXT_PRIMARY)
    services = [
        ("features-delta-one", "92s", "30s", ACCENT_AMBER),
        ("execution-service", "2s", "5s", ACCENT_GREEN),
        ("risk-and-exposure", "4s", "10s", ACCENT_GREEN),
        ("pnl-attribution", "8s", "15s", ACCENT_GREEN),
    ]
    y = 6.05
    for svc, fresh, sla, color in services:
        add_text_box(slide, 10.35, y, 1.5, 0.18, svc, font_size=7, font_name=FONT_MONO, color=TEXT_SECONDARY)
        add_text_box(slide, 11.85, y, 0.4, 0.18, fresh, font_size=7, font_name=FONT_MONO, font_color_override=color)
        add_text_box(slide, 12.3, y, 0.4, 0.18, sla, font_size=7, font_name=FONT_MONO, color=TEXT_MUTED)
        y += 0.3


def slide_02_problem(prs):
    """Before/After — real UI names, real duplication."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, 0.5, 0.4, 8, 0.4, "The core insight", font_size=28, bold=True, color=TEXT_PRIMARY)
    add_text_box(
        slide,
        0.5,
        0.85,
        10,
        0.3,
        "The dark theme and component library are already institutional. The problem is information architecture.",
        font_size=12,
        color=TEXT_SECONDARY,
    )

    # Before
    add_card(slide, 0.5, 1.5, 5.8, 4.5, BG_ELEVATED)
    add_text_box(slide, 0.7, 1.6, 5, 0.3, "Before: 11 isolated UIs", font_size=16, bold=True, color=ACCENT_RED)
    add_text_box(
        slide,
        0.7,
        1.95,
        5,
        0.3,
        "Duplicate deployments. Duplicate P&L. Duplicate alerts. Manual port-hopping.",
        font_size=9,
        color=TEXT_TERTIARY,
    )

    before_uis = [
        ("strategy-ui :5175", "execution-analytics-ui :5174", "90% identical routes"),
        ("trading-analytics-ui :5180", "client-reporting-ui :5182", "Both show P&L"),
        ("batch-audit-ui :5181", "logs-dashboard-ui :5178", "Same API :8013"),
        ("live-health-monitor-ui :5177", "settlement-ui :5176", "Both show positions"),
        ("onboarding-ui :5173", "deployment-ui :5183", "Both have /deployments"),
    ]
    y = 2.35
    for ui1, ui2, problem in before_uis:
        add_text_box(slide, 0.7, y, 2.2, 0.2, ui1, font_size=8, font_name=FONT_MONO, color=TEXT_SECONDARY)
        add_text_box(slide, 2.9, y, 0.3, 0.2, "<->", font_size=8, color=TEXT_MUTED)
        add_text_box(slide, 3.2, y, 2.2, 0.2, ui2, font_size=8, font_name=FONT_MONO, color=TEXT_SECONDARY)
        add_text_box(slide, 5.4, y, 0.8, 0.2, problem, font_size=7, font_color_override=ACCENT_RED)
        y += 0.35

    # Key failures
    failures = [
        "No workflow coherence (Design -> ... -> Reconcile invisible)",
        "No entity hierarchy (Fund -> Client -> Strategy never navigable)",
        "No cross-UI navigation (users type port numbers)",
        "/deployments in 10 of 11 UIs",
        "P&L in 3 different UIs with no canonical home",
    ]
    y = 4.2
    for f in failures:
        add_text_box(slide, 0.7, y, 5.3, 0.2, f"  {f}", font_size=9, color=TEXT_SECONDARY)
        y += 0.3

    # After
    add_card(slide, 6.8, 1.5, 6.0, 4.5, BG_ELEVATED)
    add_text_box(slide, 7.0, 1.6, 5, 0.3, "After: one operating model", font_size=16, bold=True, color=ACCENT_GREEN)
    add_text_box(
        slide,
        7.0,
        1.95,
        5,
        0.3,
        "4 daily-use surfaces + 3 specialist tools. Connected by hierarchy, lifecycle, and deep links.",
        font_size=9,
        color=TEXT_TERTIARY,
    )

    surfaces = [
        ("Trading Command Center", ":5177", "OBSERVE, INTERVENE", "Live now", ACCENT_GREEN),
        ("Strategy Analytics", ":5175", "DESIGN, SIMULATE, PROMOTE", "Historical", ACCENT_BLUE),
        ("Market Intelligence", ":5180", "EXPLAIN, RECONCILE", "Post-trade", ACCENT_PURPLE),
        ("Operations Hub", ":5183", "DEPLOY, DIAGNOSE", "Ops time", ACCENT_AMBER),
    ]
    y = 2.4
    for name, port, verbs, time_h, color in surfaces:
        add_card(slide, 7.0, y, 5.6, 0.5, BG_CARD)
        add_text_box(slide, 7.15, y + 0.05, 2.5, 0.2, name, font_size=11, bold=True, font_color_override=color)
        add_text_box(slide, 7.15, y + 0.27, 1.0, 0.18, port, font_size=7, font_name=FONT_MONO, color=TEXT_MUTED)
        add_text_box(slide, 8.3, y + 0.27, 2.0, 0.18, verbs, font_size=7, color=TEXT_TERTIARY)
        add_text_box(slide, 10.8, y + 0.15, 1.5, 0.18, time_h, font_size=8, font_color_override=color)
        y += 0.58

    tools = [
        ("Config & Onboarding", ":5173", ACCENT_CYAN),
        ("ML Platform", ":5179", ACCENT_ORANGE),
        ("Reporting & Settlement", ":5182", ACCENT_TEAL),
    ]
    y += 0.15
    add_text_box(slide, 7.0, y, 3, 0.2, "Specialist tools:", font_size=9, color=TEXT_MUTED)
    y += 0.3
    x = 7.0
    for name, port, color in tools:
        add_pill(slide, x, y, 1.7, f"{name} {port}", color, 7)
        x += 1.85

    # Thesis
    add_card(slide, 0.5, 6.3, 12.3, 0.6, BG_SECONDARY)
    add_text_box(
        slide,
        0.7,
        6.38,
        11.5,
        0.5,
        "Citadel-grade: the complexity stays; the navigation becomes legible. 7 repos, 7 ports, zero duplication.",
        font_size=12,
        color=TEXT_PRIMARY,
        bold=True,
    )


def slide_03_sma_hierarchy(prs):
    """SMA model + three parallel hierarchies (P&L, Position, Risk)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, 0.5, 0.4, 10, 0.4, "SMA model & entity hierarchy", font_size=28, bold=True, color=TEXT_PRIMARY)
    add_text_box(
        slide,
        0.5,
        0.85,
        10,
        0.3,
        "Separately managed accounts: one strategy template, isolated client instances, positions diverge by design.",
        font_size=11,
        color=TEXT_SECONDARY,
    )

    # SMA diagram
    add_card(slide, 0.5, 1.4, 6.0, 2.8, BG_ELEVATED)
    add_text_box(slide, 0.7, 1.5, 5, 0.3, "Unit of execution", font_size=14, bold=True, color=TEXT_PRIMARY)
    add_text_box(
        slide,
        0.7,
        1.85,
        5.5,
        0.2,
        "(strategy_id, client_id, config_version)  =  one async task, one position set",
        font_size=9,
        font_name=FONT_MONO,
        color=ACCENT_CYAN,
    )

    sma_rows = [
        ("DEFI_ETH_STAKED_BASIS_SCE_1H", "odum", "v2.4.0", "Live", ACCENT_GREEN),
        ("DEFI_ETH_STAKED_BASIS_SCE_1H", "apex_capital", "v2.4.0", "Live", ACCENT_GREEN),
        ("DEFI_ETH_STAKED_BASIS_SCE_1H", "meridian_fund", "v2.3.1", "Paper", ACCENT_AMBER),
        ("DEFI_ETH_BASIS_SCE_1H", "odum", "v3.2.1", "Live", ACCENT_GREEN),
        ("TRADFI_SPY_ML_DIR_1H", "odum", "v1.0.0", "Live", ACCENT_GREEN),
        ("SPORTS_HT_ML_V2.1", "apex_capital", "v2.1.0", "Live", ACCENT_GREEN),
    ]

    headers = [("Strategy Template", 0.7, 2.8), ("Client", 3.5, 0.9), ("Config", 4.5, 0.6), ("Status", 5.3, 0.5)]
    for label, x, w in headers:
        add_text_box(slide, x, 2.15, w, 0.2, label, font_size=8, bold=True, color=TEXT_MUTED)

    y = 2.4
    for strat, client, cfg, status, color in sma_rows:
        add_text_box(slide, 0.7, y, 2.8, 0.18, strat, font_size=7, font_name=FONT_MONO, color=ACCENT_CYAN)
        add_text_box(slide, 3.5, y, 0.9, 0.18, client, font_size=8, color=TEXT_SECONDARY)
        add_text_box(slide, 4.5, y, 0.6, 0.18, cfg, font_size=7, font_name=FONT_MONO, color=TEXT_TERTIARY)
        add_text_box(slide, 5.3, y, 0.5, 0.18, status, font_size=8, font_color_override=color)
        y += 0.3

    add_text_box(
        slide,
        0.7,
        3.75,
        5.5,
        0.2,
        "Same strategy code, different fills, different positions. This is expected and normal.",
        font_size=8,
        color=TEXT_MUTED,
    )

    # Three parallel hierarchies
    add_card(slide, 6.8, 1.4, 6.0, 2.8, BG_ELEVATED)
    add_text_box(
        slide,
        7.0,
        1.5,
        5,
        0.3,
        "Three parallel hierarchies (batch + live)",
        font_size=14,
        bold=True,
        color=TEXT_PRIMARY,
    )
    add_text_box(
        slide,
        7.0,
        1.85,
        5.5,
        0.2,
        "Each has a batch (historical) and live (real-time) variant. Same drill path, different data source.",
        font_size=9,
        color=TEXT_TERTIARY,
    )

    hierarchies = [
        (
            "P&L Hierarchy",
            "All -> Client -> Strategy -> Venue -> Component",
            "Funding, basis, staking yield, delta, greeks, slippage, gas, recon",
            ACCENT_GREEN,
            "Market Intelligence",
        ),
        (
            "Position Hierarchy",
            "All -> Client -> Strategy -> Venue -> Instrument",
            "Net qty, avg price, unrealized PnL, margin used, health factor",
            ACCENT_BLUE,
            "Trading Command Center (live) / Reporting (EOD)",
        ),
        (
            "Risk / Exposure Hierarchy",
            "All -> Client -> Strategy -> Venue -> Risk Dim",
            "Delta exp, funding exp, basis spread, LTV, greeks, margin util",
            ACCENT_PURPLE,
            "Trading Command Center",
        ),
    ]
    y = 2.2
    for name, path, components, color, home in hierarchies:
        add_text_box(slide, 7.0, y, 2.5, 0.2, name, font_size=10, bold=True, font_color_override=color)
        add_text_box(slide, 7.0, y + 0.22, 5.5, 0.18, path, font_size=8, font_name=FONT_MONO, color=TEXT_SECONDARY)
        add_text_box(slide, 7.0, y + 0.42, 5.5, 0.18, components, font_size=7, color=TEXT_MUTED)
        add_text_box(slide, 7.0, y + 0.6, 5.5, 0.15, f"Home: {home}", font_size=7, font_color_override=color)
        y += 0.85

    # Strategy archetypes
    add_card(slide, 0.5, 4.5, 12.3, 2.6, BG_ELEVATED)
    add_text_box(
        slide,
        0.7,
        4.6,
        10,
        0.3,
        "Strategy archetypes — groupable, not just by asset class",
        font_size=16,
        bold=True,
        color=TEXT_PRIMARY,
    )
    add_text_box(
        slide,
        0.7,
        4.95,
        10,
        0.2,
        "Strategies share patterns across asset classes. Groups by archetype, filters by asset class.",
        font_size=9,
        color=TEXT_TERTIARY,
    )

    archetypes = [
        (
            "Arbitrage",
            "Exploit price discrepancies across venues or legs",
            ["DEFI_ETH_BASIS_SCE_1H", "SPORTS_ARB", "CEFI cross-exchange"],
            ACCENT_CYAN,
        ),
        (
            "Market Making",
            "Continuous quoting, earn spread/fees",
            ["DEFI_ETH_USDT_MM_LP_V3", "CEFI_BTC_MM_BINANCE", "SPORTS_MM_BETFAIR"],
            ACCENT_AMBER,
        ),
        (
            "Directional",
            "ML-driven or signal-driven position taking",
            ["TRADFI_SPY_ML_DIR_1H", "SPORTS_HT_ML_V2.1", "CEFI_MOMENTUM"],
            ACCENT_BLUE,
        ),
        (
            "Yield",
            "Earn yield from lending, staking, or funding rates",
            ["DEFI_USDT_LENDING_SCE", "DEFI_ETH_STAKED_BASIS", "DEFI_ETH_RECURSIVE"],
            ACCENT_GREEN,
        ),
    ]

    x = 0.7
    for name, desc, examples, color in archetypes:
        add_card(slide, x, 5.3, 2.85, 1.6, BG_CARD)
        add_text_box(slide, x + 0.15, 5.4, 2.5, 0.25, name, font_size=13, bold=True, font_color_override=color)
        add_text_box(slide, x + 0.15, 5.7, 2.5, 0.3, desc, font_size=8, color=TEXT_SECONDARY)
        ey = 6.05
        for ex in examples:
            add_text_box(slide, x + 0.15, ey, 2.5, 0.18, ex, font_size=7, font_name=FONT_MONO, color=TEXT_TERTIARY)
            ey += 0.2
        x += 3.05

    add_text_box(
        slide,
        0.7,
        7.0,
        10,
        0.2,
        "ML Training is a deep-dive sub-view of Directional strategies — not a separate archetype.",
        font_size=9,
        color=TEXT_MUTED,
        bold=True,
    )


def slide_04_exploration_to_live(prs):
    """Exploration-to-Live config flow — shared pattern across Strategy, ML, Execution."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(
        slide,
        0.5,
        0.4,
        10,
        0.4,
        "Exploration to live — one pattern, three domain views",
        font_size=28,
        bold=True,
        color=TEXT_PRIMARY,
    )
    add_text_box(
        slide,
        0.5,
        0.85,
        10,
        0.3,
        "Strategy, ML, and Execution configs follow the same flow shape. Each domain adds its own nuances.",
        font_size=11,
        color=TEXT_SECONDARY,
    )

    # Shared flow steps
    add_card(slide, 0.5, 1.4, 12.3, 1.0, BG_ELEVATED)
    add_text_box(
        slide, 0.7, 1.5, 10, 0.2, "Shared flow shape (all three domains):", font_size=10, bold=True, color=TEXT_PRIMARY
    )

    steps = [
        "Configure\nGrid Params",
        "Generate\nConfigs",
        "Run Batch\n(Sharded)",
        "Analyse\nResults",
        "Select\nCandidates",
        "Promote\nto Live",
    ]
    x = 0.7
    for i, step in enumerate(steps):
        c = ACCENT_CYAN if i < 3 else ACCENT_GREEN if i < 5 else ACCENT_AMBER
        add_card(slide, x, 1.8, 1.8, 0.45, BG_CARD)
        lines = step.split("\n")
        add_text_box(slide, x + 0.05, 1.82, 1.7, 0.2, lines[0], font_size=9, bold=True, font_color_override=c)
        if len(lines) > 1:
            add_text_box(slide, x + 0.05, 2.0, 1.7, 0.18, lines[1], font_size=8, color=TEXT_TERTIARY)
        if i < 5:
            add_text_box(slide, x + 1.85, 1.92, 0.15, 0.2, "->", font_size=10, color=TEXT_MUTED)
        x += 2.0

    # Three domain columns
    domains = [
        (
            "Strategy Config",
            ACCENT_BLUE,
            "Strategy Analytics",
            [
                "Instruments + venues + algos",
                "param grid: mode x timeframe x category",
                "Sharded by category/venue/date",
                "Sharpe, PnL, drawdown, win rate",
                "DimensionalGrid multi-select",
                "Cross-link to Ops /deploy",
            ],
        ),
        (
            "ML Model Config",
            ACCENT_ORANGE,
            "ML Platform",
            [
                "Model arch + hyperparameters",
                "grid: lr x layers x dropout x timeframe",
                "GPU training jobs (Cloud Run)",
                "Accuracy, loss, Sharpe improvement",
                "Best model per instrument",
                "Deploy to ml-inference-service",
            ],
        ),
        (
            "Execution Config",
            ACCENT_PURPLE,
            "Strategy Analytics (Execution tab)",
            [
                "Algo type + venue routing",
                "grid: algo x instrument x benchmark",
                "Replay on historical tick data",
                "Alpha bps, slippage, fill rate",
                "Best algo per venue+instrument",
                "Update execution-service config",
            ],
        ),
    ]

    x = 0.5
    for name, color, surface, steps in domains:
        add_card(slide, x, 2.6, 3.9, 4.2, BG_ELEVATED)
        add_text_box(slide, x + 0.2, 2.7, 3.5, 0.25, name, font_size=14, bold=True, font_color_override=color)
        add_text_box(slide, x + 0.2, 3.0, 3.5, 0.2, f"Surface: {surface}", font_size=8, color=TEXT_MUTED)

        labels = ["Configure", "Generate", "Run Batch", "Analyse", "Select", "Promote"]
        y = 3.3
        for i, (label, detail) in enumerate(zip(labels, steps)):
            add_text_box(
                slide, x + 0.2, y, 0.8, 0.18, f"0{i + 1} {label}", font_size=8, bold=True, font_color_override=color
            )
            add_text_box(slide, x + 1.1, y, 2.5, 0.18, detail, font_size=8, color=TEXT_SECONDARY)
            y += 0.38
        x += 4.1

    add_card(slide, 0.5, 7.0, 12.3, 0.4, BG_SECONDARY)
    add_text_box(
        slide,
        0.7,
        7.05,
        11.5,
        0.3,
        "No auto-deploy. Promote creates reviewed handoff via cross-link to Ops Hub with deploy form.",
        font_size=10,
        color=TEXT_PRIMARY,
        bold=True,
    )


def slide_05_command_center(prs):
    """Trading Command Center — full detail with real data."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_nav_bar(slide, active_idx=0, show_as_of=True, as_of_text="Live")
    add_lifecycle_rail(slide, active_idx=3, top=0.48)
    add_breadcrumb(slide, ["Odum Delta One", "All Clients", "All Strategies"], top=0.9)

    add_text_box(
        slide, 0.3, 1.1, 8, 0.3, "Surface 01 — Trading Command Center", font_size=20, bold=True, color=ACCENT_GREEN
    )
    add_text_box(
        slide,
        0.3,
        1.45,
        8,
        0.25,
        "What is the state of the world — right now, or at any point in time?",
        font_size=10,
        color=TEXT_SECONDARY,
    )

    # KPI row with sparklines
    add_kpi_card(slide, 0.3, 1.85, "Firm PnL", "+$1.42m", "  +0.8% 1d", ACCENT_GREEN, "up")
    add_kpi_card(slide, 2.5, 1.85, "Net Exposure", "$4.2m", "1.2x levered", ACCENT_CYAN, "volatile")
    add_kpi_card(slide, 4.7, 1.85, "Margin Health", "82%", "$340k free", ACCENT_AMBER, "flat")
    add_kpi_card(slide, 6.9, 1.85, "Live Strats", "12", "9 healthy", ACCENT_BLUE, "flat")
    add_kpi_card(slide, 9.1, 1.85, "Alerts", "3 crit", "2 high", ACCENT_RED, "down")

    # Kill switch panel
    add_card(slide, 0.3, 3.15, 12.7, 1.1, BG_ELEVATED)
    add_text_box(slide, 0.5, 3.2, 3, 0.25, "Kill Switch / Intervention", font_size=12, bold=True, color=ACCENT_RED)
    add_text_box(
        slide,
        0.5,
        3.5,
        12,
        0.2,
        "Scope:  Fund > Client > Strategy > Venue   |   "
        "Actions:  Pause strategy  |  Cancel orders  |  Flatten exposure  |  Disable venue",
        font_size=9,
        color=TEXT_SECONDARY,
    )
    add_text_box(
        slide,
        0.5,
        3.8,
        12,
        0.2,
        "Every action requires rationale, shows position count, generates incident + audit record.",
        font_size=8,
        color=TEXT_MUTED,
    )

    # Routes
    add_card(slide, 0.3, 4.45, 6.2, 2.8, BG_ELEVATED)
    add_text_box(slide, 0.5, 4.55, 5, 0.25, "Routes", font_size=12, bold=True, color=TEXT_PRIMARY)
    routes = [
        ("/", "Fund dashboard: KPIs + strategy table + P&L/risk attribution + alerts"),
        ("/positions", "Live positions, filter by [Client | Strategy | Venue | Asset Class]"),
        ("/positions/:runId", "Position detail: orders, fills, execution timeline"),
        ("/risk", "Risk matrix + exposure attribution + margin health + DeFi LTV"),
        ("/alerts", "Unified alert feed, severity-colored, incident creation"),
        ("/health", "Service health grid + dependency DAG + feature freshness SLA"),
        ("/manual", "Manual trade entry (scoped) + kill switches"),
    ]
    y = 4.85
    for route, desc in routes:
        add_text_box(slide, 0.5, y, 1.5, 0.18, route, font_size=8, font_name=FONT_MONO, color=ACCENT_CYAN)
        add_text_box(slide, 2.0, y, 4.3, 0.18, desc, font_size=8, color=TEXT_SECONDARY)
        y += 0.33

    # Risk tabs detail
    add_card(slide, 6.7, 4.45, 6.3, 2.8, BG_ELEVATED)
    add_text_box(slide, 6.9, 4.55, 5, 0.25, "/risk — Four tabs", font_size=12, bold=True, color=TEXT_PRIMARY)
    risk_tabs = [
        ("Risk Summary", "All limits at a glance, highest utilization first. Value vs limit with %.", ACCENT_BLUE),
        ("Exposure Attribution", "Same 6D as P&L, each with limit threshold. Forward-looking.", ACCENT_PURPLE),
        ("Margin & LTV", "Per-venue margin vs limit. Per-position DeFi health factor vs liquidation.", ACCENT_AMBER),
        (
            "Limits Detail",
            "Full hierarchy drill-down: firm -> client -> strategy -> venue -> instrument.",
            ACCENT_GREEN,
        ),
    ]
    y = 4.9
    for tab, desc, color in risk_tabs:
        add_text_box(slide, 6.9, y, 2.0, 0.2, tab, font_size=10, bold=True, font_color_override=color)
        add_text_box(slide, 6.9, y + 0.22, 5.8, 0.3, desc, font_size=8, color=TEXT_SECONDARY)
        y += 0.55

    # Limit bars example
    add_text_box(
        slide,
        6.9,
        7.0,
        5.8,
        0.15,
        "Delta Exp  ████████████░░░░░  $2.4m/$5.0m (48%) ●  |  Margin  █████████████████░  78%/80% (97%) ▲",
        font_size=7,
        font_name=FONT_MONO,
        color=TEXT_TERTIARY,
    )


def slide_06_strategy_analytics(prs):
    """Strategy Analytics — DimensionalGrid with real data."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_nav_bar(slide, active_idx=1)
    add_lifecycle_rail(slide, active_idx=2, top=0.48)
    add_breadcrumb(slide, ["Strategy Analytics", "Grid Results", "Selected Candidates"], top=0.9)

    add_text_box(slide, 0.3, 1.1, 8, 0.3, "Surface 02 — Strategy Analytics", font_size=20, bold=True, color=ACCENT_BLUE)
    add_text_box(
        slide,
        0.3,
        1.45,
        8,
        0.25,
        "Catalogue -> Backtests -> DimensionalGrid -> Select best -> Promote to live. The quant's home.",
        font_size=10,
        color=TEXT_SECONDARY,
    )

    # Catalogue strip
    add_card(slide, 0.3, 1.85, 3.5, 1.4, BG_ELEVATED)
    add_text_box(slide, 0.5, 1.95, 3, 0.2, "Strategy Catalogue", font_size=11, bold=True, color=TEXT_PRIMARY)
    add_text_box(slide, 0.5, 2.2, 3, 0.15, "Group by archetype, filter by asset class", font_size=8, color=TEXT_MUTED)

    archetypes = [
        ("Yield (4)", ACCENT_GREEN),
        ("Arbitrage (2)", ACCENT_CYAN),
        ("Directional (3)", ACCENT_BLUE),
        ("Market Making (3)", ACCENT_AMBER),
    ]
    y = 2.45
    for arch, color in archetypes:
        add_text_box(slide, 0.5, y, 1.5, 0.18, arch, font_size=9, font_color_override=color)
        y += 0.22

    # DimensionalGrid
    add_card(slide, 4.0, 1.85, 9.0, 4.0, BG_ELEVATED)
    add_text_box(
        slide,
        4.2,
        1.95,
        5,
        0.25,
        "DimensionalGrid — the killer feature for batch analysis",
        font_size=12,
        bold=True,
        color=TEXT_PRIMARY,
    )

    # Dimension pills
    dims = [("Instrument", 4.2), ("Venue", 5.5), ("Strategy", 6.6), ("Date", 7.6), ("Config", 8.5)]
    for label, x in dims:
        add_pill(slide, x, 2.25, 1.0, label, ACCENT_CYAN, 7)

    add_text_box(slide, 10.2, 2.28, 1.0, 0.2, "Showing 47 of 1,203", font_size=7, font_name=FONT_MONO, color=TEXT_MUTED)
    add_pill(slide, 11.4, 2.25, 0.8, "Heatmap", ACCENT_PURPLE, 7)
    add_pill(slide, 12.3, 2.25, 0.6, "Export", TEXT_SECONDARY, 7)

    # Grid headers
    grid_cols = [
        (" ", 4.2, 0.25),
        ("Experiment", 4.5, 1.0),
        ("Strategy", 5.5, 1.5),
        ("Config", 7.0, 0.8),
        ("Venue", 7.8, 0.8),
        ("Shard", 8.6, 0.7),
        ("Sharpe", 9.3, 0.6),
        ("PnL", 9.9, 0.7),
        ("DD", 10.6, 0.5),
        ("Status", 11.2, 0.7),
    ]
    for label, x, w in grid_cols:
        add_text_box(slide, x, 2.6, w, 0.18, label, font_size=7, bold=True, color=TEXT_MUTED)

    # Grid rows — real data
    grid_rows = [
        (True, "exp-221", "DEFI_ETH_BASIS", "3.3.0-rc1", "Bin/Hyper", "2025Q4", "2.1", "+$1.8m", "4.1%", "selected"),
        (True, "exp-301", "DEFI_ETH_STAKED", "2.5.0", "EtherFi/HL", "2026Q1", "2.5", "+$2.4m", "3.3%", "selected"),
        (False, "exp-222", "DEFI_ETH_BASIS", "3.3.0-rc2", "Bin/OKX", "2026Q1", "1.7", "+$1.1m", "5.1%", "review"),
        (False, "exp-711", "TRADFI_SPY_ML", "5.1.2", "IBKR", "2025H2", "1.3", "+$0.7m", "6.8%", "hold"),
        (True, "exp-402", "SPORTS_HT_ML", "2.1.0", "Betfair", "2026Q1", "1.6", "+$0.3m", "1.8%", "selected"),
    ]

    y = 2.85
    for selected, exp, strat, cfg, venue, shard, sharpe, pnl, dd, status in grid_rows:
        check = "[x]" if selected else "[ ]"
        add_text_box(
            slide,
            4.2,
            y,
            0.25,
            0.18,
            check,
            font_size=8,
            font_name=FONT_MONO,
            font_color_override=ACCENT_GREEN if selected else TEXT_MUTED,
        )
        add_text_box(slide, 4.5, y, 1.0, 0.18, exp, font_size=8, font_name=FONT_MONO, color=TEXT_SECONDARY)
        add_text_box(slide, 5.5, y, 1.5, 0.18, strat, font_size=8, font_name=FONT_MONO, color=ACCENT_CYAN)
        add_text_box(slide, 7.0, y, 0.8, 0.18, cfg, font_size=7, font_name=FONT_MONO, color=TEXT_TERTIARY)
        add_text_box(slide, 7.8, y, 0.8, 0.18, venue, font_size=8, color=TEXT_SECONDARY)
        add_text_box(slide, 8.6, y, 0.7, 0.18, shard, font_size=8, color=TEXT_TERTIARY)
        add_text_box(slide, 9.3, y, 0.6, 0.18, sharpe, font_size=8, font_name=FONT_MONO, color=TEXT_PRIMARY)
        add_text_box(slide, 9.9, y, 0.7, 0.18, pnl, font_size=8, font_name=FONT_MONO, font_color_override=ACCENT_GREEN)
        add_text_box(slide, 10.6, y, 0.5, 0.18, dd, font_size=8, font_name=FONT_MONO, color=TEXT_SECONDARY)
        status_color = ACCENT_GREEN if status == "selected" else ACCENT_AMBER if status == "review" else TEXT_MUTED
        add_text_box(slide, 11.2, y, 0.7, 0.18, status, font_size=7, font_color_override=status_color)
        y += 0.38

    # Selection toolbar
    add_card(slide, 4.2, 4.85, 8.6, 0.4, BG_SECONDARY)
    add_text_box(
        slide,
        4.4,
        4.9,
        8,
        0.3,
        "Selection: 3 configs   |   [Promote to Batch]   [Promote to Live]   [Export CSV]",
        font_size=9,
        bold=True,
        color=TEXT_PRIMARY,
    )

    # Promotion package
    add_card(slide, 0.3, 3.5, 3.5, 2.3, BG_ELEVATED)
    add_text_box(slide, 0.5, 3.6, 3, 0.25, "Promotion Package", font_size=11, bold=True, color=ACCENT_AMBER)
    promo_items = [
        "Source batch run + shard",
        "Capacity estimate per instrument",
        "Max drawdown observed",
        "Regime notes (trending/mean-rev)",
        "Target environment (staging/live)",
        "Approval status + owner",
    ]
    y = 3.9
    for item in promo_items:
        add_text_box(slide, 0.5, y, 3, 0.18, f"  {item}", font_size=8, color=TEXT_SECONDARY)
        y += 0.22

    add_text_box(slide, 0.5, 5.35, 3, 0.2, "Cross-link to Ops -> /deploy", font_size=9, bold=True, color=ACCENT_AMBER)

    # Routes
    add_card(slide, 0.3, 5.9, 12.7, 1.3, BG_ELEVATED)
    add_text_box(slide, 0.5, 6.0, 3, 0.2, "Key routes", font_size=11, bold=True, color=TEXT_PRIMARY)

    route_cols = [
        [
            ("/strategies", "Catalogue: filter by archetype + asset class"),
            ("/strategies/:id", "Hub: Overview | Live | Backtest | Results | Execution | Deep Dive"),
            ("/grid", "DimensionalGrid: slice by [strategy, instrument, venue, date, config]"),
        ],
        [
            ("/compare", "Overlay equity curves, rank by metric"),
            ("/heatmap", "Two-dimension color matrix (instrument x algo -> Sharpe)"),
            ("/generate", "Config generator -> mass deploy -> progress tracking"),
        ],
        [
            ("/configs", "Browse all configs across strategies"),
            ("/instruments", "Instrument definitions + data availability"),
            ("/tick-data", "Market tick data explorer with candle/tick toggle"),
        ],
    ]
    x = 0.5
    for col in route_cols:
        y = 6.25
        for route, desc in col:
            add_text_box(slide, x, y, 1.3, 0.18, route, font_size=7, font_name=FONT_MONO, color=ACCENT_CYAN)
            add_text_box(slide, x + 1.35, y, 2.7, 0.18, desc, font_size=7, color=TEXT_SECONDARY)
            y += 0.28
        x += 4.2


def slide_07_markets_ops(prs):
    """Market Intelligence + Operations Hub side by side."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(
        slide, 0.5, 0.4, 10, 0.4, "Surfaces 03-04 — Markets + Ops", font_size=28, bold=True, color=TEXT_PRIMARY
    )
    add_text_box(
        slide,
        0.5,
        0.85,
        10,
        0.3,
        "Separate explanation from deployment. Connect them through entity links and correlation IDs.",
        font_size=11,
        color=TEXT_SECONDARY,
    )

    # Market Intelligence
    add_card(slide, 0.5, 1.4, 6.2, 5.5, BG_ELEVATED)
    add_text_box(
        slide, 0.7, 1.5, 5, 0.3, "Market Intelligence", font_size=16, bold=True, font_color_override=ACCENT_PURPLE
    )
    add_text_box(
        slide,
        0.7,
        1.85,
        5,
        0.2,
        "Explain P&L. Drill 5 levels. Reconcile residuals. Batch + live variants.",
        font_size=9,
        color=TEXT_TERTIARY,
    )

    # P&L waterfall — real components
    add_text_box(
        slide,
        0.7,
        2.2,
        3,
        0.2,
        "P&L Waterfall (real components per strategy type)",
        font_size=10,
        bold=True,
        color=TEXT_PRIMARY,
    )

    waterfall = [
        ("Funding Rate", "+$412k", 0.85),
        ("Basis Convergence", "+$355k", 0.72),
        ("Staking Yield (weETH)", "+$145k", 0.40),
        ("Delta PnL", "+$61k", 0.18),
        ("Rewards (EIGEN)", "+$22k", 0.08),
        ("Execution Slippage", "-$61k", -0.18),
        ("Gas + Swap Fees", "-$44k", -0.13),
        ("Recon Pending", "-$18k", -0.06),
    ]
    y = 2.5
    for comp, val, _bar in waterfall:
        add_text_box(slide, 0.7, y, 1.8, 0.18, comp, font_size=8, color=TEXT_SECONDARY)
        c = ACCENT_GREEN if val.startswith("+") else ACCENT_RED
        add_text_box(slide, 2.6, y, 0.8, 0.18, val, font_size=8, font_name=FONT_MONO, font_color_override=c)
        # Mini bar
        bar_w = abs(_bar) * 3.0
        bar_x = 3.5 if _bar >= 0 else 3.5 - bar_w
        bar_color = ACCENT_GREEN if _bar >= 0 else ACCENT_RED
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bar_x), Inches(y + 0.02), Inches(bar_w), Inches(0.14)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bar_color
        shape.line.fill.background()
        y += 0.25

    # Drill path
    add_text_box(
        slide,
        0.7,
        4.65,
        5,
        0.2,
        "Drill path: All -> Client -> Strategy -> Venue -> Component",
        font_size=8,
        font_name=FONT_MONO,
        color=ACCENT_PURPLE,
    )

    # Recon
    add_text_box(slide, 0.7, 5.0, 3, 0.2, "Batch vs Live Reconciliation", font_size=10, bold=True, color=TEXT_PRIMARY)
    recon_items = [
        "What backtest predicted vs what live executed",
        "ML strategy signals vs actual fills",
        "Execution service domain data comparison",
        "Residual bucket links to open recon cases",
        "Correlation ID traces across services",
    ]
    y = 5.3
    for item in recon_items:
        add_text_box(slide, 0.7, y, 5.5, 0.18, f"  {item}", font_size=8, color=TEXT_SECONDARY)
        y += 0.22

    # Routes
    add_text_box(
        slide,
        0.7,
        6.3,
        5,
        0.2,
        "Routes: /pnl, /pnl/client/:id, /pnl/strategy/:id,",
        font_size=7,
        font_name=FONT_MONO,
        color=TEXT_MUTED,
    )
    add_text_box(
        slide,
        0.7,
        6.5,
        5,
        0.2,
        "/desk, /orderbook, /latency, /recon, /reports",
        font_size=7,
        font_name=FONT_MONO,
        color=TEXT_MUTED,
    )

    # Operations Hub
    add_card(slide, 7.0, 1.4, 5.8, 5.5, BG_ELEVATED)
    add_text_box(slide, 7.2, 1.5, 5, 0.3, "Operations Hub", font_size=16, bold=True, font_color_override=ACCENT_AMBER)
    add_text_box(
        slide,
        7.2,
        1.85,
        5,
        0.2,
        "Deploy, observe, diagnose. Never pretend to be the P&L home.",
        font_size=9,
        color=TEXT_TERTIARY,
    )

    # Batch summary
    add_card(slide, 7.2, 2.2, 5.3, 1.2, BG_CARD)
    add_text_box(
        slide, 7.4, 2.3, 3, 0.2, "Batch Summary + Data Completeness", font_size=10, bold=True, color=TEXT_PRIMARY
    )

    add_text_box(
        slide,
        7.4,
        2.55,
        2,
        0.2,
        "47 done  |  3 failed  |  12 running",
        font_size=9,
        font_name=FONT_MONO,
        color=TEXT_SECONDARY,
    )

    services_deploy = [
        ("execution-service", "prod-tokyo v31"),
        ("features-delta-one", "prod-eu v17"),
        ("risk-and-exposure", "prod-eu v12"),
    ]
    y = 2.85
    for svc, deploy in services_deploy:
        add_text_box(slide, 7.4, y, 2.5, 0.18, svc, font_size=8, font_name=FONT_MONO, color=ACCENT_CYAN)
        add_text_box(slide, 9.9, y, 2.5, 0.18, deploy, font_size=8, color=TEXT_TERTIARY)
        y += 0.22

    # Sidebar nav groups
    add_text_box(slide, 7.2, 3.6, 5, 0.2, "Sidebar navigation groups", font_size=10, bold=True, color=TEXT_PRIMARY)

    groups = [
        ("DEPLOY", ["Overview", "Deploy Service", "Services", "Epics"], ACCENT_AMBER),
        ("OBSERVE", ["Batch Jobs", "Logs", "Events", "Data Health"], ACCENT_BLUE),
        ("COMPLIANCE", ["Audit Trail", "Compliance", "CI/CD"], ACCENT_PURPLE),
    ]
    y = 3.9
    for group, items, color in groups:
        add_text_box(slide, 7.2, y, 1.5, 0.2, group, font_size=9, bold=True, font_color_override=color)
        add_text_box(slide, 8.5, y, 4, 0.2, "  |  ".join(items), font_size=8, color=TEXT_SECONDARY)
        y += 0.35

    # Event hierarchy
    add_card(slide, 7.2, 4.8, 5.3, 0.8, BG_CARD)
    add_text_box(
        slide, 7.4, 4.85, 5, 0.2, "Event -> Alert -> Incident hierarchy", font_size=9, bold=True, color=ACCENT_BLUE
    )
    add_text_box(
        slide,
        7.4,
        5.1,
        5,
        0.45,
        "Raw event (log_event()) -> lands in GCS events/{svc}/{date}/events.jsonl\n"
        "If severity >= threshold -> Alert (Trading CC or Ops, by type)\n"
        "If escalated -> Incident (kill switch, manual, auto-rules)",
        font_size=7,
        font_name=FONT_MONO,
        color=TEXT_SECONDARY,
    )

    # Logging & correlation
    add_card(slide, 7.2, 5.7, 2.5, 0.7, BG_CARD)
    add_text_box(slide, 7.35, 5.75, 2.3, 0.2, "Logs & correlation", font_size=9, bold=True, color=ACCENT_CYAN)
    add_text_box(
        slide,
        7.35,
        6.0,
        2.3,
        0.35,
        "Filter: service, severity,\ntime, text, correlation_id\nClick corr_id -> full trace",
        font_size=7,
        color=TEXT_SECONDARY,
    )

    # Version alignment
    add_card(slide, 9.85, 5.7, 2.65, 0.7, BG_CARD)
    add_text_box(slide, 10.0, 5.75, 2.4, 0.2, "Version alignment", font_size=9, bold=True, color=ACCENT_AMBER)
    add_text_box(
        slide,
        10.0,
        6.0,
        2.4,
        0.35,
        "/services grid shows:\ndeployed vs expected version\ndrift = amber badge",
        font_size=7,
        color=TEXT_SECONDARY,
    )

    # Grafana + shared services
    add_card(slide, 7.2, 6.55, 5.3, 0.55, BG_CARD)
    add_text_box(
        slide, 7.4, 6.58, 2.5, 0.18, "Grafana: SRE deep-dive only", font_size=8, bold=True, color=ACCENT_PURPLE
    )
    add_text_box(slide, 10.0, 6.58, 2.3, 0.18, "All services are shared", font_size=8, bold=True, color=TEXT_SECONDARY)
    add_text_box(
        slide,
        7.4,
        6.78,
        5,
        0.25,
        "Service detail -> 'Open in Grafana' (new tab, pre-filtered)  |  "
        "Sharding is logical (strategy_id, client_id), not physical",
        font_size=7,
        color=TEXT_MUTED,
    )


def slide_08_implementation(prs):
    """Implementation path + outcomes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, 0.5, 0.4, 10, 0.4, "Implementation path", font_size=28, bold=True, color=TEXT_PRIMARY)
    add_text_box(
        slide,
        0.5,
        0.85,
        10,
        0.3,
        "Build the shared language first, then collapse the duplicated surfaces behind it.",
        font_size=11,
        color=TEXT_SECONDARY,
    )

    # Phase stack
    phases = [
        (
            "P0",
            "Shared ui-kit",
            "GlobalNavBar, LifecycleRail, BreadcrumbNav, EntityLink, CrossLink, DimensionalGrid, FilterBar",
            ACCENT_CYAN,
        ),
        (
            "P1",
            "Remove /deployments",
            "Mount nav shell everywhere. Replace plain-text entities with deep links.",
            ACCENT_BLUE,
        ),
        (
            "P2-P5",
            "Collapse duplicate UIs",
            "Strategy merge (exec-analytics + strategy). Ops merge (batch-audit + logs + deploy). Settlement.",
            ACCENT_PURPLE,
        ),
        (
            "P6",
            "Canonical P&L hierarchy",
            "5-level drill-down (All->Client->Strategy->Venue->Component). 6D attribution. Batch vs live recon.",
            ACCENT_GREEN,
        ),
        (
            "P7-P8",
            "Hardening + config flows",
            "FilterBars + breadcrumbs across all surfaces. Accessibility. Real config generator wiring. E2E tests.",
            ACCENT_AMBER,
        ),
    ]

    y = 1.3
    for phase, title, desc, color in phases:
        add_card(slide, 0.5, y, 12.3, 0.65, BG_ELEVATED)
        add_text_box(slide, 0.7, y + 0.05, 0.6, 0.25, phase, font_size=14, bold=True, font_color_override=color)
        add_text_box(slide, 1.4, y + 0.05, 2.5, 0.25, title, font_size=12, bold=True, color=TEXT_PRIMARY)
        add_text_box(slide, 1.4, y + 0.32, 11, 0.25, desc, font_size=9, color=TEXT_SECONDARY)
        y += 0.75

    # Foundational artifacts
    add_card(slide, 0.5, 5.3, 6.0, 1.8, BG_ELEVATED)
    add_text_box(
        slide,
        0.7,
        5.4,
        5,
        0.25,
        "Foundational artifacts (already created)",
        font_size=12,
        bold=True,
        color=TEXT_PRIMARY,
    )

    artifacts = [
        ("canonical-surface-ownership.mdc", "Ownership table + overlap resolution"),
        ("component-patterns.mdc", "Design tokens, AppShell, EntityLink patterns"),
        ("cross-surface-navigation.mdc", "Entity routing map + URL param protocol"),
        ("ui-quality-gates.mdc", "8-gate standard: lint, types, vitest, build, a11y"),
        ("dimensional-grid-spec.mdc", "Props, UX behavior, promotion flow, perf reqs"),
    ]
    y = 5.7
    for file, desc in artifacts:
        add_text_box(slide, 0.7, y, 2.5, 0.18, file, font_size=7, font_name=FONT_MONO, color=ACCENT_CYAN)
        add_text_box(slide, 3.3, y, 3, 0.18, desc, font_size=8, color=TEXT_SECONDARY)
        y += 0.27

    # Outcomes
    add_card(slide, 6.8, 5.3, 6.0, 1.8, BG_ELEVATED)
    add_text_box(slide, 7.0, 5.4, 5, 0.25, "What 'done' looks like", font_size=12, bold=True, color=TEXT_PRIMARY)

    outcomes = [
        (
            "User",
            [
                "One answer to 'what is happening right now?' — at any point in time",
                "One canonical P&L home with 5-level drill-down + time series",
                "Risk limits shown as value vs threshold, highest util first",
            ],
        ),
        (
            "System",
            [
                "150+ endpoints, 0 orphans — every API output has a UI home",
                "7 repos, 7 ports, zero duplicated routes",
                "Sports/DeFi/CeFi/TradFi via FilterBar dimension, not separate surfaces",
            ],
        ),
    ]
    y = 5.7
    for category, items in outcomes:
        add_text_box(slide, 7.0, y, 1.0, 0.18, category, font_size=9, bold=True, color=ACCENT_GREEN)
        for item in items:
            add_text_box(slide, 8.0, y, 4.5, 0.18, item, font_size=8, color=TEXT_SECONDARY)
            y += 0.25

    # Final thesis
    add_card(slide, 0.5, 7.2, 12.3, 0.25, BG_SECONDARY)
    add_text_box(
        slide,
        0.7,
        7.2,
        11.5,
        0.25,
        "The complexity stays. The navigation becomes legible: hierarchy, lifecycle, cross-linking.",
        font_size=11,
        bold=True,
        color=TEXT_PRIMARY,
        alignment=PP_ALIGN.CENTER,
    )


def slide_09_temporal(prs):
    """Temporal universality — snapshot + time series everywhere."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(
        slide,
        0.5,
        0.4,
        10,
        0.4,
        "Temporal universality — snapshot + time series everywhere",
        font_size=28,
        bold=True,
        color=TEXT_PRIMARY,
    )
    add_text_box(
        slide,
        0.5,
        0.85,
        10,
        0.3,
        "Every view showing state supports two temporal modes. This is an overarching design principle.",
        font_size=11,
        color=TEXT_SECONDARY,
    )

    # Left: Point-in-time snapshot
    add_card(slide, 0.5, 1.4, 6.0, 3.0, BG_ELEVATED)
    add_text_box(
        slide,
        0.7,
        1.5,
        5,
        0.3,
        "Point-in-time snapshot (the time machine)",
        font_size=14,
        bold=True,
        color=ACCENT_AMBER,
    )
    add_text_box(
        slide,
        0.7,
        1.85,
        5.5,
        0.4,
        "Pick any datetime. See the exact state of the world at that moment. Same layout, same drill-down, "
        "just at a different point in time.",
        font_size=9,
        color=TEXT_SECONDARY,
    )

    # AsOfDatetimePicker mockup
    add_card(slide, 0.7, 2.4, 5.5, 0.5, BG_CARD)
    add_text_box(slide, 0.9, 2.45, 1.5, 0.2, "AsOfDatetimePicker", font_size=8, bold=True, color=TEXT_MUTED)
    options = [
        ("Live", ACCENT_GREEN),
        ("1h ago", TEXT_SECONDARY),
        ("Start of day", TEXT_SECONDARY),
        ("Yesterday close", TEXT_SECONDARY),
        ("2026-03-14 14:32", ACCENT_AMBER),
    ]
    x = 0.9
    for label, color in options:
        add_pill(slide, x, 2.65, 1.2 if len(label) > 8 else 0.8, label, color, 7)
        x += 1.05 if len(label) > 8 else 0.85

    add_text_box(
        slide,
        0.7,
        3.05,
        5.5,
        0.2,
        "URL:  ?as_of=2026-03-14T14:32:00Z   (omitted = live, shareable, bookmarkable)",
        font_size=8,
        font_name=FONT_MONO,
        color=TEXT_TERTIARY,
    )

    # Two screenshots side by side: Live vs Historical
    add_card(slide, 0.7, 3.35, 2.6, 0.8, BG_CARD)
    add_pill(slide, 0.8, 3.4, 0.6, "Live", ACCENT_GREEN, 7)
    add_text_box(slide, 0.8, 3.65, 2.3, 0.2, "Firm PnL", font_size=8, color=TEXT_TERTIARY)
    add_text_box(
        slide,
        0.8,
        3.85,
        2.3,
        0.2,
        "+$1.42m",
        font_size=14,
        font_name=FONT_MONO,
        bold=True,
        font_color_override=ACCENT_GREEN,
    )

    add_card(slide, 3.5, 3.35, 2.6, 0.8, BG_CARD)
    add_pill(slide, 3.6, 3.4, 1.3, "2026-03-14 14:32", ACCENT_AMBER, 7)
    add_text_box(slide, 3.6, 3.65, 2.3, 0.2, "Firm PnL", font_size=8, color=TEXT_TERTIARY)
    add_text_box(
        slide,
        3.6,
        3.85,
        2.3,
        0.2,
        "+$0.89m",
        font_size=14,
        font_name=FONT_MONO,
        bold=True,
        font_color_override=ACCENT_AMBER,
    )

    # Right: Time series
    add_card(slide, 6.8, 1.4, 6.0, 3.0, BG_ELEVATED)
    add_text_box(
        slide, 7.0, 1.5, 5, 0.3, "Time series (trajectory embedded inline)", font_size=14, bold=True, color=ACCENT_BLUE
    )
    add_text_box(
        slide,
        7.0,
        1.85,
        5.5,
        0.4,
        "Every KPI card shows a sparkline. Every table metric has a 'show over time' toggle. "
        "The trajectory matters as much as the value.",
        font_size=9,
        color=TEXT_SECONDARY,
    )

    # KPI with sparkline example
    add_card(slide, 7.0, 2.4, 2.5, 1.4, BG_CARD)
    add_text_box(slide, 7.15, 2.5, 2, 0.15, "Firm PnL", font_size=9, color=TEXT_TERTIARY)
    add_text_box(
        slide,
        7.15,
        2.7,
        2,
        0.3,
        "+$1.42m",
        font_size=20,
        font_name=FONT_MONO,
        bold=True,
        font_color_override=ACCENT_GREEN,
    )
    add_text_box(slide, 7.15, 3.0, 2, 0.15, "+0.8% 1d", font_size=8, font_name=FONT_MONO, color=TEXT_MUTED)
    add_sparkline(slide, 7.15, 3.2, 2.1, 0.45, ACCENT_GREEN, "up")
    add_text_box(slide, 7.15, 3.65, 2, 0.12, "1d  1w  1m  3m  1y", font_size=7, color=TEXT_MUTED)

    # TimeSeriesToggle example
    add_card(slide, 9.7, 2.4, 2.8, 1.4, BG_CARD)
    add_text_box(slide, 9.85, 2.5, 2.5, 0.15, "Delta Exposure", font_size=9, color=TEXT_TERTIARY)
    add_text_box(
        slide,
        9.85,
        2.7,
        2.5,
        0.3,
        "$2.4m",
        font_size=20,
        font_name=FONT_MONO,
        bold=True,
        font_color_override=ACCENT_CYAN,
    )
    add_text_box(slide, 9.85, 3.0, 2.5, 0.15, "net long", font_size=8, font_name=FONT_MONO, color=TEXT_MUTED)
    add_sparkline(slide, 9.85, 3.2, 2.3, 0.45, ACCENT_CYAN, "volatile")
    add_text_box(slide, 9.85, 3.65, 2.5, 0.12, "1d  1w  1m  3m  1y", font_size=7, color=TEXT_MUTED)

    # What gets snapshot + time series table
    add_card(slide, 0.5, 4.6, 12.3, 2.8, BG_ELEVATED)
    add_text_box(
        slide, 0.7, 4.7, 10, 0.25, "What gets snapshot + time series", font_size=14, bold=True, color=TEXT_PRIMARY
    )

    domains = [
        (
            "Positions",
            "Exact positions at T",
            "Position size evolution, entry/exit markers",
            "Trading Command Center",
            ACCENT_GREEN,
        ),
        (
            "P&L",
            "Cumulative P&L at T, attribution at T",
            "Equity curve, daily bars, component evolution",
            "Market Intelligence",
            ACCENT_PURPLE,
        ),
        (
            "Risk / Exposure",
            "Delta/funding/basis/greeks at T",
            "Exposure drift, margin trend, LTV evolution",
            "Trading Command Center",
            ACCENT_CYAN,
        ),
        (
            "Alerts",
            "Which alerts were active at T",
            "Alert frequency, resolution time trends",
            "Trading Command Center",
            ACCENT_RED,
        ),
        (
            "Strategy Status",
            "Which strategies live/paused at T",
            "Lifecycle timeline (promoted, paused, killed)",
            "Strategy Analytics",
            ACCENT_BLUE,
        ),
        (
            "Feature Freshness",
            "SLA compliance at T",
            "Freshness degradation, breach history",
            "Trading Command Center",
            ACCENT_AMBER,
        ),
        (
            "Orders / Fills",
            "Orders in flight at T",
            "Order flow over time, fill rate",
            "Market Intelligence",
            ACCENT_PURPLE,
        ),
    ]

    headers = [
        ("Domain", 0.7, 1.5),
        ("Snapshot (as-of T)", 2.2, 2.5),
        ("Time Series (over period)", 4.8, 3.0),
        ("Canonical Home", 8.0, 2.0),
    ]
    for label, x, w in headers:
        add_text_box(slide, x, 5.0, w, 0.2, label, font_size=8, bold=True, color=TEXT_MUTED)

    y = 5.25
    for domain, snapshot, timeseries, home, color in domains:
        add_text_box(slide, 0.7, y, 1.5, 0.18, domain, font_size=9, bold=True, font_color_override=color)
        add_text_box(slide, 2.2, y, 2.5, 0.18, snapshot, font_size=8, color=TEXT_SECONDARY)
        add_text_box(slide, 4.8, y, 3.0, 0.18, timeseries, font_size=8, color=TEXT_SECONDARY)
        add_text_box(slide, 8.0, y, 2.0, 0.18, home, font_size=8, font_color_override=color)
        y += 0.3


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════


def main():
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W))
    prs.slide_height = Emu(int(SLIDE_H))

    slide_01_hero(prs)
    slide_02_problem(prs)
    slide_03_sma_hierarchy(prs)
    slide_04_exploration_to_live(prs)
    slide_05_command_center(prs)
    slide_06_strategy_analytics(prs)
    slide_07_markets_ops(prs)
    slide_09_temporal(prs)
    slide_08_implementation(prs)

    output = "unified-trading-pm/plans/active/UI_Platform_Redesign_Vision_v2.pptx"
    prs.save(output)
    print(f"Saved: {output}")
    print("  9 slides with temporal universality, sparklines, as-of picker")


if __name__ == "__main__":
    main()
