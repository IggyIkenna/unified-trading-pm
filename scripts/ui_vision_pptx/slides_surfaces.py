"""Slides 05-07 — Trading Command Center, Strategy Analytics, Markets + Ops."""

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from .deck_style import (
    ACCENT_AMBER,
    ACCENT_BLUE,
    ACCENT_CYAN,
    ACCENT_GREEN,
    ACCENT_PURPLE,
    ACCENT_RED,
    BG_CARD,
    BG_ELEVATED,
    BG_SECONDARY,
    FONT_MONO,
    TEXT_MUTED,
    TEXT_PRIMARY,
    TEXT_SECONDARY,
    TEXT_TERTIARY,
    add_breadcrumb,
    add_card,
    add_kpi_card,
    add_lifecycle_rail,
    add_nav_bar,
    add_pill,
    add_text_box,
    set_slide_bg,
)


def slide_05_command_center(prs):
    """Trading Command Center — full detail with real data."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_nav_bar(slide, active_idx=0, show_as_of=True, as_of_text="Live")
    add_lifecycle_rail(slide, active_idx=3, top=0.48)
    add_breadcrumb(slide, ["Odum Delta One", "All Clients", "All Strategies"], top=0.9)

    add_text_box(
        slide, 0.3, 1.1, 8, 0.3, "Surface 01 — Trading Command Center", font_size=20, bold=True, color=ACCENT_GREEN
    )
    add_text_box(
        slide,
        0.3,
        1.45,
        8,
        0.25,
        "What is the state of the world — right now, or at any point in time?",
        font_size=10,
        color=TEXT_SECONDARY,
    )

    # KPI row with sparklines
    add_kpi_card(slide, 0.3, 1.85, "Firm PnL", "+$1.42m", "  +0.8% 1d", ACCENT_GREEN, "up")
    add_kpi_card(slide, 2.5, 1.85, "Net Exposure", "$4.2m", "1.2x levered", ACCENT_CYAN, "volatile")
    add_kpi_card(slide, 4.7, 1.85, "Margin Health", "82%", "$340k free", ACCENT_AMBER, "flat")
    add_kpi_card(slide, 6.9, 1.85, "Live Strats", "12", "9 healthy", ACCENT_BLUE, "flat")
    add_kpi_card(slide, 9.1, 1.85, "Alerts", "3 crit", "2 high", ACCENT_RED, "down")

    # Kill switch panel
    add_card(slide, 0.3, 3.15, 12.7, 1.1, BG_ELEVATED)
    add_text_box(slide, 0.5, 3.2, 3, 0.25, "Kill Switch / Intervention", font_size=12, bold=True, color=ACCENT_RED)
    add_text_box(
        slide,
        0.5,
        3.5,
        12,
        0.2,
        "Scope:  Fund > Client > Strategy > Venue   |   "
        "Actions:  Pause strategy  |  Cancel orders  |  Flatten exposure  |  Disable venue",
        font_size=9,
        color=TEXT_SECONDARY,
    )
    add_text_box(
        slide,
        0.5,
        3.8,
        12,
        0.2,
        "Every action requires rationale, shows position count, generates incident + audit record.",
        font_size=8,
        color=TEXT_MUTED,
    )

    # Routes
    add_card(slide, 0.3, 4.45, 6.2, 2.8, BG_ELEVATED)
    add_text_box(slide, 0.5, 4.55, 5, 0.25, "Routes", font_size=12, bold=True, color=TEXT_PRIMARY)
    routes = [
        ("/", "Fund dashboard: KPIs + strategy table + P&L/risk attribution + alerts"),
        ("/positions", "Live positions, filter by [Client | Strategy | Venue | Asset Class]"),
        ("/positions/:runId", "Position detail: orders, fills, execution timeline"),
        ("/risk", "Risk matrix + exposure attribution + margin health + DeFi LTV"),
        ("/alerts", "Unified alert feed, severity-colored, incident creation"),
        ("/health", "Service health grid + dependency DAG + feature freshness SLA"),
        ("/manual", "Manual trade entry (scoped) + kill switches"),
    ]
    y = 4.85
    for route, desc in routes:
        add_text_box(slide, 0.5, y, 1.5, 0.18, route, font_size=8, font_name=FONT_MONO, color=ACCENT_CYAN)
        add_text_box(slide, 2.0, y, 4.3, 0.18, desc, font_size=8, color=TEXT_SECONDARY)
        y += 0.33

    # Risk tabs detail
    add_card(slide, 6.7, 4.45, 6.3, 2.8, BG_ELEVATED)
    add_text_box(slide, 6.9, 4.55, 5, 0.25, "/risk — Four tabs", font_size=12, bold=True, color=TEXT_PRIMARY)
    risk_tabs = [
        ("Risk Summary", "All limits at a glance, highest utilization first. Value vs limit with %.", ACCENT_BLUE),
        ("Exposure Attribution", "Same 6D as P&L, each with limit threshold. Forward-looking.", ACCENT_PURPLE),
        ("Margin & LTV", "Per-venue margin vs limit. Per-position DeFi health factor vs liquidation.", ACCENT_AMBER),
        (
            "Limits Detail",
            "Full hierarchy drill-down: firm -> client -> strategy -> venue -> instrument.",
            ACCENT_GREEN,
        ),
    ]
    y = 4.9
    for tab, desc, color in risk_tabs:
        add_text_box(slide, 6.9, y, 2.0, 0.2, tab, font_size=10, bold=True, font_color_override=color)
        add_text_box(slide, 6.9, y + 0.22, 5.8, 0.3, desc, font_size=8, color=TEXT_SECONDARY)
        y += 0.55

    # Limit bars example
    add_text_box(
        slide,
        6.9,
        7.0,
        5.8,
        0.15,
        "Delta Exp  ████████████░░░░░  $2.4m/$5.0m (48%) ●  |  Margin  █████████████████░  78%/80% (97%) ▲",
        font_size=7,
        font_name=FONT_MONO,
        color=TEXT_TERTIARY,
    )


def slide_06_strategy_analytics(prs):
    """Strategy Analytics — DimensionalGrid with real data."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_nav_bar(slide, active_idx=1)
    add_lifecycle_rail(slide, active_idx=2, top=0.48)
    add_breadcrumb(slide, ["Strategy Analytics", "Grid Results", "Selected Candidates"], top=0.9)

    add_text_box(slide, 0.3, 1.1, 8, 0.3, "Surface 02 — Strategy Analytics", font_size=20, bold=True, color=ACCENT_BLUE)
    add_text_box(
        slide,
        0.3,
        1.45,
        8,
        0.25,
        "Catalogue -> Backtests -> DimensionalGrid -> Select best -> Promote to live. The quant's home.",
        font_size=10,
        color=TEXT_SECONDARY,
    )

    # Catalogue strip
    add_card(slide, 0.3, 1.85, 3.5, 1.4, BG_ELEVATED)
    add_text_box(slide, 0.5, 1.95, 3, 0.2, "Strategy Catalogue", font_size=11, bold=True, color=TEXT_PRIMARY)
    add_text_box(slide, 0.5, 2.2, 3, 0.15, "Group by archetype, filter by asset class", font_size=8, color=TEXT_MUTED)

    archetypes = [
        ("Yield (4)", ACCENT_GREEN),
        ("Arbitrage (2)", ACCENT_CYAN),
        ("Directional (3)", ACCENT_BLUE),
        ("Market Making (3)", ACCENT_AMBER),
    ]
    y = 2.45
    for arch, color in archetypes:
        add_text_box(slide, 0.5, y, 1.5, 0.18, arch, font_size=9, font_color_override=color)
        y += 0.22

    # DimensionalGrid
    add_card(slide, 4.0, 1.85, 9.0, 4.0, BG_ELEVATED)
    add_text_box(
        slide,
        4.2,
        1.95,
        5,
        0.25,
        "DimensionalGrid — the killer feature for batch analysis",
        font_size=12,
        bold=True,
        color=TEXT_PRIMARY,
    )

    # Dimension pills
    dims = [("Instrument", 4.2), ("Venue", 5.5), ("Strategy", 6.6), ("Date", 7.6), ("Config", 8.5)]
    for label, x in dims:
        add_pill(slide, x, 2.25, 1.0, label, ACCENT_CYAN, 7)

    add_text_box(slide, 10.2, 2.28, 1.0, 0.2, "Showing 47 of 1,203", font_size=7, font_name=FONT_MONO, color=TEXT_MUTED)
    add_pill(slide, 11.4, 2.25, 0.8, "Heatmap", ACCENT_PURPLE, 7)
    add_pill(slide, 12.3, 2.25, 0.6, "Export", TEXT_SECONDARY, 7)

    # Grid headers
    grid_cols = [
        (" ", 4.2, 0.25),
        ("Experiment", 4.5, 1.0),
        ("Strategy", 5.5, 1.5),
        ("Config", 7.0, 0.8),
        ("Venue", 7.8, 0.8),
        ("Shard", 8.6, 0.7),
        ("Sharpe", 9.3, 0.6),
        ("PnL", 9.9, 0.7),
        ("DD", 10.6, 0.5),
        ("Status", 11.2, 0.7),
    ]
    for label, x, w in grid_cols:
        add_text_box(slide, x, 2.6, w, 0.18, label, font_size=7, bold=True, color=TEXT_MUTED)

    # Grid rows — real data
    grid_rows = [
        (True, "exp-221", "DEFI_ETH_BASIS", "3.3.0-rc1", "Bin/Hyper", "2025Q4", "2.1", "+$1.8m", "4.1%", "selected"),
        (True, "exp-301", "DEFI_ETH_STAKED", "2.5.0", "EtherFi/HL", "2026Q1", "2.5", "+$2.4m", "3.3%", "selected"),
        (False, "exp-222", "DEFI_ETH_BASIS", "3.3.0-rc2", "Bin/OKX", "2026Q1", "1.7", "+$1.1m", "5.1%", "review"),
        (False, "exp-711", "TRADFI_SPY_ML", "5.1.2", "IBKR", "2025H2", "1.3", "+$0.7m", "6.8%", "hold"),
        (True, "exp-402", "SPORTS_HT_ML", "2.1.0", "Betfair", "2026Q1", "1.6", "+$0.3m", "1.8%", "selected"),
    ]

    y = 2.85
    for selected, exp, strat, cfg, venue, shard, sharpe, pnl, dd, status in grid_rows:
        check = "[x]" if selected else "[ ]"
        add_text_box(
            slide,
            4.2,
            y,
            0.25,
            0.18,
            check,
            font_size=8,
            font_name=FONT_MONO,
            font_color_override=ACCENT_GREEN if selected else TEXT_MUTED,
        )
        add_text_box(slide, 4.5, y, 1.0, 0.18, exp, font_size=8, font_name=FONT_MONO, color=TEXT_SECONDARY)
        add_text_box(slide, 5.5, y, 1.5, 0.18, strat, font_size=8, font_name=FONT_MONO, color=ACCENT_CYAN)
        add_text_box(slide, 7.0, y, 0.8, 0.18, cfg, font_size=7, font_name=FONT_MONO, color=TEXT_TERTIARY)
        add_text_box(slide, 7.8, y, 0.8, 0.18, venue, font_size=8, color=TEXT_SECONDARY)
        add_text_box(slide, 8.6, y, 0.7, 0.18, shard, font_size=8, color=TEXT_TERTIARY)
        add_text_box(slide, 9.3, y, 0.6, 0.18, sharpe, font_size=8, font_name=FONT_MONO, color=TEXT_PRIMARY)
        add_text_box(slide, 9.9, y, 0.7, 0.18, pnl, font_size=8, font_name=FONT_MONO, font_color_override=ACCENT_GREEN)
        add_text_box(slide, 10.6, y, 0.5, 0.18, dd, font_size=8, font_name=FONT_MONO, color=TEXT_SECONDARY)
        status_color = ACCENT_GREEN if status == "selected" else ACCENT_AMBER if status == "review" else TEXT_MUTED
        add_text_box(slide, 11.2, y, 0.7, 0.18, status, font_size=7, font_color_override=status_color)
        y += 0.38

    # Selection toolbar
    add_card(slide, 4.2, 4.85, 8.6, 0.4, BG_SECONDARY)
    add_text_box(
        slide,
        4.4,
        4.9,
        8,
        0.3,
        "Selection: 3 configs   |   [Promote to Batch]   [Promote to Live]   [Export CSV]",
        font_size=9,
        bold=True,
        color=TEXT_PRIMARY,
    )

    # Promotion package
    add_card(slide, 0.3, 3.5, 3.5, 2.3, BG_ELEVATED)
    add_text_box(slide, 0.5, 3.6, 3, 0.25, "Promotion Package", font_size=11, bold=True, color=ACCENT_AMBER)
    promo_items = [
        "Source batch run + shard",
        "Capacity estimate per instrument",
        "Max drawdown observed",
        "Regime notes (trending/mean-rev)",
        "Target environment (staging/live)",
        "Approval status + owner",
    ]
    y = 3.9
    for item in promo_items:
        add_text_box(slide, 0.5, y, 3, 0.18, f"  {item}", font_size=8, color=TEXT_SECONDARY)
        y += 0.22

    add_text_box(slide, 0.5, 5.35, 3, 0.2, "Cross-link to Ops -> /deploy", font_size=9, bold=True, color=ACCENT_AMBER)

    # Routes
    add_card(slide, 0.3, 5.9, 12.7, 1.3, BG_ELEVATED)
    add_text_box(slide, 0.5, 6.0, 3, 0.2, "Key routes", font_size=11, bold=True, color=TEXT_PRIMARY)

    route_cols = [
        [
            ("/strategies", "Catalogue: filter by archetype + asset class"),
            ("/strategies/:id", "Hub: Overview | Live | Backtest | Results | Execution | Deep Dive"),
            ("/grid", "DimensionalGrid: slice by [strategy, instrument, venue, date, config]"),
        ],
        [
            ("/compare", "Overlay equity curves, rank by metric"),
            ("/heatmap", "Two-dimension color matrix (instrument x algo -> Sharpe)"),
            ("/generate", "Config generator -> mass deploy -> progress tracking"),
        ],
        [
            ("/configs", "Browse all configs across strategies"),
            ("/instruments", "Instrument definitions + data availability"),
            ("/tick-data", "Market tick data explorer with candle/tick toggle"),
        ],
    ]
    x = 0.5
    for col in route_cols:
        y = 6.25
        for route, desc in col:
            add_text_box(slide, x, y, 1.3, 0.18, route, font_size=7, font_name=FONT_MONO, color=ACCENT_CYAN)
            add_text_box(slide, x + 1.35, y, 2.7, 0.18, desc, font_size=7, color=TEXT_SECONDARY)
            y += 0.28
        x += 4.2


def slide_07_markets_ops(prs):
    """Market Intelligence + Operations Hub side by side."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(
        slide, 0.5, 0.4, 10, 0.4, "Surfaces 03-04 — Markets + Ops", font_size=28, bold=True, color=TEXT_PRIMARY
    )
    add_text_box(
        slide,
        0.5,
        0.85,
        10,
        0.3,
        "Separate explanation from deployment. Connect them through entity links and correlation IDs.",
        font_size=11,
        color=TEXT_SECONDARY,
    )

    # Market Intelligence
    add_card(slide, 0.5, 1.4, 6.2, 5.5, BG_ELEVATED)
    add_text_box(
        slide, 0.7, 1.5, 5, 0.3, "Market Intelligence", font_size=16, bold=True, font_color_override=ACCENT_PURPLE
    )
    add_text_box(
        slide,
        0.7,
        1.85,
        5,
        0.2,
        "Explain P&L. Drill 5 levels. Reconcile residuals. Batch + live variants.",
        font_size=9,
        color=TEXT_TERTIARY,
    )

    # P&L waterfall — real components
    add_text_box(
        slide,
        0.7,
        2.2,
        3,
        0.2,
        "P&L Waterfall (real components per strategy type)",
        font_size=10,
        bold=True,
        color=TEXT_PRIMARY,
    )

    waterfall = [
        ("Funding Rate", "+$412k", 0.85),
        ("Basis Convergence", "+$355k", 0.72),
        ("Staking Yield (weETH)", "+$145k", 0.40),
        ("Delta PnL", "+$61k", 0.18),
        ("Rewards (EIGEN)", "+$22k", 0.08),
        ("Execution Slippage", "-$61k", -0.18),
        ("Gas + Swap Fees", "-$44k", -0.13),
        ("Recon Pending", "-$18k", -0.06),
    ]
    y = 2.5
    for comp, val, _bar in waterfall:
        add_text_box(slide, 0.7, y, 1.8, 0.18, comp, font_size=8, color=TEXT_SECONDARY)
        c = ACCENT_GREEN if val.startswith("+") else ACCENT_RED
        add_text_box(slide, 2.6, y, 0.8, 0.18, val, font_size=8, font_name=FONT_MONO, font_color_override=c)
        # Mini bar
        bar_w = abs(_bar) * 3.0
        bar_x = 3.5 if _bar >= 0 else 3.5 - bar_w
        bar_color = ACCENT_GREEN if _bar >= 0 else ACCENT_RED
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bar_x), Inches(y + 0.02), Inches(bar_w), Inches(0.14)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bar_color
        shape.line.fill.background()
        y += 0.25

    # Drill path
    add_text_box(
        slide,
        0.7,
        4.65,
        5,
        0.2,
        "Drill path: All -> Client -> Strategy -> Venue -> Component",
        font_size=8,
        font_name=FONT_MONO,
        color=ACCENT_PURPLE,
    )

    # Recon
    add_text_box(slide, 0.7, 5.0, 3, 0.2, "Batch vs Live Reconciliation", font_size=10, bold=True, color=TEXT_PRIMARY)
    recon_items = [
        "What backtest predicted vs what live executed",
        "ML strategy signals vs actual fills",
        "Execution service domain data comparison",
        "Residual bucket links to open recon cases",
        "Correlation ID traces across services",
    ]
    y = 5.3
    for item in recon_items:
        add_text_box(slide, 0.7, y, 5.5, 0.18, f"  {item}", font_size=8, color=TEXT_SECONDARY)
        y += 0.22

    # Routes
    add_text_box(
        slide,
        0.7,
        6.3,
        5,
        0.2,
        "Routes: /pnl, /pnl/client/:id, /pnl/strategy/:id,",
        font_size=7,
        font_name=FONT_MONO,
        color=TEXT_MUTED,
    )
    add_text_box(
        slide,
        0.7,
        6.5,
        5,
        0.2,
        "/desk, /orderbook, /latency, /recon, /reports",
        font_size=7,
        font_name=FONT_MONO,
        color=TEXT_MUTED,
    )

    # Operations Hub
    add_card(slide, 7.0, 1.4, 5.8, 5.5, BG_ELEVATED)
    add_text_box(slide, 7.2, 1.5, 5, 0.3, "Operations Hub", font_size=16, bold=True, font_color_override=ACCENT_AMBER)
    add_text_box(
        slide,
        7.2,
        1.85,
        5,
        0.2,
        "Deploy, observe, diagnose. Never pretend to be the P&L home.",
        font_size=9,
        color=TEXT_TERTIARY,
    )

    # Batch summary
    add_card(slide, 7.2, 2.2, 5.3, 1.2, BG_CARD)
    add_text_box(
        slide, 7.4, 2.3, 3, 0.2, "Batch Summary + Data Completeness", font_size=10, bold=True, color=TEXT_PRIMARY
    )

    add_text_box(
        slide,
        7.4,
        2.55,
        2,
        0.2,
        "47 done  |  3 failed  |  12 running",
        font_size=9,
        font_name=FONT_MONO,
        color=TEXT_SECONDARY,
    )

    services_deploy = [
        ("execution-service", "prod-tokyo v31"),
        ("features-delta-one", "prod-eu v17"),
        ("risk-and-exposure", "prod-eu v12"),
    ]
    y = 2.85
    for svc, deploy in services_deploy:
        add_text_box(slide, 7.4, y, 2.5, 0.18, svc, font_size=8, font_name=FONT_MONO, color=ACCENT_CYAN)
        add_text_box(slide, 9.9, y, 2.5, 0.18, deploy, font_size=8, color=TEXT_TERTIARY)
        y += 0.22

    # Sidebar nav groups
    add_text_box(slide, 7.2, 3.6, 5, 0.2, "Sidebar navigation groups", font_size=10, bold=True, color=TEXT_PRIMARY)

    groups = [
        ("DEPLOY", ["Overview", "Deploy Service", "Services", "Epics"], ACCENT_AMBER),
        ("OBSERVE", ["Batch Jobs", "Logs", "Events", "Data Health"], ACCENT_BLUE),
        ("COMPLIANCE", ["Audit Trail", "Compliance", "CI/CD"], ACCENT_PURPLE),
    ]
    y = 3.9
    for group, items, color in groups:
        add_text_box(slide, 7.2, y, 1.5, 0.2, group, font_size=9, bold=True, font_color_override=color)
        add_text_box(slide, 8.5, y, 4, 0.2, "  |  ".join(items), font_size=8, color=TEXT_SECONDARY)
        y += 0.35

    # Event hierarchy
    add_card(slide, 7.2, 4.8, 5.3, 0.8, BG_CARD)
    add_text_box(
        slide, 7.4, 4.85, 5, 0.2, "Event -> Alert -> Incident hierarchy", font_size=9, bold=True, color=ACCENT_BLUE
    )
    add_text_box(
        slide,
        7.4,
        5.1,
        5,
        0.45,
        "Raw event (log_event()) -> lands in GCS events/{svc}/{date}/events.jsonl\n"
        "If severity >= threshold -> Alert (Trading CC or Ops, by type)\n"
        "If escalated -> Incident (kill switch, manual, auto-rules)",
        font_size=7,
        font_name=FONT_MONO,
        color=TEXT_SECONDARY,
    )

    # Logging & correlation
    add_card(slide, 7.2, 5.7, 2.5, 0.7, BG_CARD)
    add_text_box(slide, 7.35, 5.75, 2.3, 0.2, "Logs & correlation", font_size=9, bold=True, color=ACCENT_CYAN)
    add_text_box(
        slide,
        7.35,
        6.0,
        2.3,
        0.35,
        "Filter: service, severity,\ntime, text, correlation_id\nClick corr_id -> full trace",
        font_size=7,
        color=TEXT_SECONDARY,
    )

    # Version alignment
    add_card(slide, 9.85, 5.7, 2.65, 0.7, BG_CARD)
    add_text_box(slide, 10.0, 5.75, 2.4, 0.2, "Version alignment", font_size=9, bold=True, color=ACCENT_AMBER)
    add_text_box(
        slide,
        10.0,
        6.0,
        2.4,
        0.35,
        "/services grid shows:\ndeployed vs expected version\ndrift = amber badge",
        font_size=7,
        color=TEXT_SECONDARY,
    )

    # Grafana + shared services
    add_card(slide, 7.2, 6.55, 5.3, 0.55, BG_CARD)
    add_text_box(
        slide, 7.4, 6.58, 2.5, 0.18, "Grafana: SRE deep-dive only", font_size=8, bold=True, color=ACCENT_PURPLE
    )
    add_text_box(slide, 10.0, 6.58, 2.3, 0.18, "All services are shared", font_size=8, bold=True, color=TEXT_SECONDARY)
    add_text_box(
        slide,
        7.4,
        6.78,
        5,
        0.25,
        "Service detail -> 'Open in Grafana' (new tab, pre-filtered)  |  "
        "Sharding is logical (strategy_id, client_id), not physical",
        font_size=7,
        color=TEXT_MUTED,
    )
