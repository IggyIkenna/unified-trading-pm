"""Design tokens + shared drawing primitives for the UI vision deck."""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Design tokens ──────────────────────────────────────────────────────────────

BG_PRIMARY = RGBColor(0x0A, 0x0A, 0x0B)
BG_SECONDARY = RGBColor(0x11, 0x11, 0x13)
BG_TERTIARY = RGBColor(0x18, 0x18, 0x1B)
BG_ELEVATED = RGBColor(0x1C, 0x1C, 0x1F)
BG_CARD = RGBColor(0x1E, 0x1E, 0x22)

TEXT_PRIMARY = RGBColor(0xFA, 0xFA, 0xFA)
TEXT_SECONDARY = RGBColor(0xA1, 0xA1, 0xAA)
TEXT_TERTIARY = RGBColor(0x71, 0x71, 0x7A)
TEXT_MUTED = RGBColor(0x52, 0x52, 0x5B)

ACCENT_CYAN = RGBColor(0x22, 0xD3, 0xEE)
ACCENT_GREEN = RGBColor(0x4A, 0xDE, 0x80)
ACCENT_AMBER = RGBColor(0xFB, 0xBF, 0x24)
ACCENT_RED = RGBColor(0xF8, 0x71, 0x71)
ACCENT_BLUE = RGBColor(0x60, 0xA5, 0xFA)
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
ACCENT_TEAL = RGBColor(0x2D, 0xD4, 0xBF)
ACCENT_ORANGE = RGBColor(0xF9, 0x73, 0x16)

BORDER = RGBColor(0x27, 0x27, 0x2A)
BORDER_SUBTLE = RGBColor(0x1F, 0x1F, 0x23)

FONT_SANS = "IBM Plex Sans"
FONT_MONO = "JetBrains Mono"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color=BG_PRIMARY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=12,
    font_name=FONT_SANS,
    color=TEXT_PRIMARY,
    bold=False,
    alignment=PP_ALIGN.LEFT,
    font_color_override=None,
):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))  # noqa: N806
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.color.rgb = font_color_override or color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_card(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Subtle border
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.5)
    shape.rotation = 0.0
    return shape


def add_pill(slide, left, top, width, text, color=ACCENT_CYAN, font_size=9):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.28))
    shape.fill.solid()
    # Dim background version of the color
    min(color[0] + 10, 255)
    min(color[1] + 10, 255)
    min(color[2] + 10, 255)
    shape.fill.fore_color.rgb = RGBColor(max(color[0] // 5, 15), max(color[1] // 5, 15), max(color[2] // 5, 15))
    shape.line.color.rgb = RGBColor(
        min(color[0] // 2 + 30, 200), min(color[1] // 2 + 30, 200), min(color[2] // 2 + 30, 200)
    )
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = FONT_SANS
    p.font.color.rgb = color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_sparkline(slide, left, top, width, height, color, pattern="up"):
    """Draw a simple sparkline shape using connected line segments."""
    if pattern == "up":
        points = [(0, 0.7), (0.15, 0.6), (0.3, 0.65), (0.45, 0.4), (0.6, 0.5), (0.75, 0.3), (0.9, 0.15), (1.0, 0.1)]
    elif pattern == "down":
        points = [(0, 0.2), (0.15, 0.3), (0.3, 0.25), (0.45, 0.5), (0.6, 0.4), (0.75, 0.6), (0.9, 0.7), (1.0, 0.8)]
    elif pattern == "flat":
        points = [(0, 0.5), (0.15, 0.45), (0.3, 0.55), (0.45, 0.5), (0.6, 0.48), (0.75, 0.52), (0.9, 0.5), (1.0, 0.49)]
    else:  # volatile
        points = [(0, 0.5), (0.15, 0.3), (0.3, 0.7), (0.45, 0.2), (0.6, 0.6), (0.75, 0.35), (0.9, 0.55), (1.0, 0.4)]
    # Draw as small dots (shapes can't do polylines easily, use tiny rectangles)
    for px, py in points:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(left + px * width), Inches(top + py * height), Inches(0.04), Inches(0.04)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()


def add_kpi_card(slide, left, top, label, value, sub, accent=ACCENT_CYAN, sparkline_pattern="up"):
    add_card(slide, left, top, 2.0, 1.3, BG_ELEVATED)
    add_text_box(slide, left + 0.15, top + 0.1, 1.7, 0.25, label, font_size=9, color=TEXT_TERTIARY)
    add_text_box(
        slide,
        left + 0.15,
        top + 0.35,
        1.7,
        0.35,
        value,
        font_size=22,
        font_name=FONT_MONO,
        bold=True,
        font_color_override=accent,
    )
    add_text_box(slide, left + 0.15, top + 0.78, 1.7, 0.25, sub, font_size=8, color=TEXT_MUTED, font_name=FONT_MONO)
    # Sparkline under the value
    add_sparkline(slide, left + 0.15, top + 0.95, 1.6, 0.25, accent, sparkline_pattern)


def add_nav_bar(slide, active_idx=0, show_as_of=False, as_of_text="Live"):
    """Add the global nav bar at the top."""
    add_card(slide, 0, 0, 13.333, 0.42, BG_SECONDARY)
    labels = [
        ("Trading", ACCENT_GREEN),
        ("Strategy", ACCENT_BLUE),
        ("Markets", ACCENT_PURPLE),
        ("Ops", ACCENT_AMBER),
        ("Config", ACCENT_CYAN),
        ("ML", ACCENT_ORANGE),
        ("Reports", ACCENT_TEAL),
    ]
    x = 0.3
    for i, (label, color) in enumerate(labels):
        c = color if i == active_idx else TEXT_MUTED
        add_text_box(slide, x, 0.08, 1.0, 0.3, label, font_size=10, bold=(i == active_idx), font_color_override=c)
        x += 1.15
    if show_as_of:
        add_pill(
            slide, 8.8, 0.08, 1.5, f"As-of: {as_of_text}", ACCENT_AMBER if as_of_text != "Live" else ACCENT_GREEN, 8
        )
    add_text_box(slide, 10.5, 0.08, 2.5, 0.3, "Search entity, run, service...", font_size=9, color=TEXT_MUTED)


def add_lifecycle_rail(slide, active_idx=3, top=0.48):
    """Add the lifecycle rail below nav."""
    steps = ["Design", "Simulate", "Promote", "Run", "Monitor", "Explain", "Reconcile"]
    x = 0.3
    for i, step in enumerate(steps):
        if i == active_idx:
            add_card(slide, x, top, 1.55, 0.35, RGBColor(0xFA, 0xFA, 0xFA))
            add_text_box(
                slide, x + 0.05, top + 0.02, 0.3, 0.15, f"0{i + 1}", font_size=7, font_color_override=BG_PRIMARY
            )
            add_text_box(
                slide, x + 0.05, top + 0.12, 1.4, 0.2, step, font_size=11, bold=True, font_color_override=BG_PRIMARY
            )
        else:
            add_card(slide, x, top, 1.55, 0.35, BG_ELEVATED)
            add_text_box(slide, x + 0.05, top + 0.02, 0.3, 0.15, f"0{i + 1}", font_size=7, color=TEXT_MUTED)
            add_text_box(slide, x + 0.05, top + 0.12, 1.4, 0.2, step, font_size=11, color=TEXT_SECONDARY)
        x += 1.72


def add_breadcrumb(slide, parts, top=0.9):
    text = "  >  ".join(parts)
    add_text_box(slide, 0.3, top, 10, 0.25, text, font_size=9, font_name=FONT_MONO, color=TEXT_TERTIARY)
