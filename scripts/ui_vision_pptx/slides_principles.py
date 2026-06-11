"""Slides 08-09 — implementation path + temporal universality."""

from pptx.enum.text import PP_ALIGN

from .deck_style import (
    ACCENT_AMBER,
    ACCENT_BLUE,
    ACCENT_CYAN,
    ACCENT_GREEN,
    ACCENT_PURPLE,
    ACCENT_RED,
    BG_CARD,
    BG_ELEVATED,
    BG_SECONDARY,
    FONT_MONO,
    TEXT_MUTED,
    TEXT_PRIMARY,
    TEXT_SECONDARY,
    TEXT_TERTIARY,
    add_card,
    add_pill,
    add_sparkline,
    add_text_box,
    set_slide_bg,
)


def slide_08_implementation(prs):
    """Implementation path + outcomes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, 0.5, 0.4, 10, 0.4, "Implementation path", font_size=28, bold=True, color=TEXT_PRIMARY)
    add_text_box(
        slide,
        0.5,
        0.85,
        10,
        0.3,
        "Build the shared language first, then collapse the duplicated surfaces behind it.",
        font_size=11,
        color=TEXT_SECONDARY,
    )

    # Phase stack
    phases = [
        (
            "P0",
            "Shared ui-kit",
            "GlobalNavBar, LifecycleRail, BreadcrumbNav, EntityLink, CrossLink, DimensionalGrid, FilterBar",
            ACCENT_CYAN,
        ),
        (
            "P1",
            "Remove /deployments",
            "Mount nav shell everywhere. Replace plain-text entities with deep links.",
            ACCENT_BLUE,
        ),
        (
            "P2-P5",
            "Collapse duplicate UIs",
            "Strategy merge (exec-analytics + strategy). Ops merge (batch-audit + logs + deploy). Settlement.",
            ACCENT_PURPLE,
        ),
        (
            "P6",
            "Canonical P&L hierarchy",
            "5-level drill-down (All->Client->Strategy->Venue->Component). 6D attribution. Batch vs live recon.",
            ACCENT_GREEN,
        ),
        (
            "P7-P8",
            "Hardening + config flows",
            "FilterBars + breadcrumbs across all surfaces. Accessibility. Real config generator wiring. E2E tests.",
            ACCENT_AMBER,
        ),
    ]

    y = 1.3
    for phase, title, desc, color in phases:
        add_card(slide, 0.5, y, 12.3, 0.65, BG_ELEVATED)
        add_text_box(slide, 0.7, y + 0.05, 0.6, 0.25, phase, font_size=14, bold=True, font_color_override=color)
        add_text_box(slide, 1.4, y + 0.05, 2.5, 0.25, title, font_size=12, bold=True, color=TEXT_PRIMARY)
        add_text_box(slide, 1.4, y + 0.32, 11, 0.25, desc, font_size=9, color=TEXT_SECONDARY)
        y += 0.75

    # Foundational artifacts
    add_card(slide, 0.5, 5.3, 6.0, 1.8, BG_ELEVATED)
    add_text_box(
        slide,
        0.7,
        5.4,
        5,
        0.25,
        "Foundational artifacts (already created)",
        font_size=12,
        bold=True,
        color=TEXT_PRIMARY,
    )

    artifacts = [
        ("canonical-surface-ownership.mdc", "Ownership table + overlap resolution"),
        ("component-patterns.mdc", "Design tokens, AppShell, EntityLink patterns"),
        ("cross-surface-navigation.mdc", "Entity routing map + URL param protocol"),
        ("ui-quality-gates.mdc", "8-gate standard: lint, types, vitest, build, a11y"),
        ("dimensional-grid-spec.mdc", "Props, UX behavior, promotion flow, perf reqs"),
    ]
    y = 5.7
    for file, desc in artifacts:
        add_text_box(slide, 0.7, y, 2.5, 0.18, file, font_size=7, font_name=FONT_MONO, color=ACCENT_CYAN)
        add_text_box(slide, 3.3, y, 3, 0.18, desc, font_size=8, color=TEXT_SECONDARY)
        y += 0.27

    # Outcomes
    add_card(slide, 6.8, 5.3, 6.0, 1.8, BG_ELEVATED)
    add_text_box(slide, 7.0, 5.4, 5, 0.25, "What 'done' looks like", font_size=12, bold=True, color=TEXT_PRIMARY)

    outcomes = [
        (
            "User",
            [
                "One answer to 'what is happening right now?' — at any point in time",
                "One canonical P&L home with 5-level drill-down + time series",
                "Risk limits shown as value vs threshold, highest util first",
            ],
        ),
        (
            "System",
            [
                "150+ endpoints, 0 orphans — every API output has a UI home",
                "7 repos, 7 ports, zero duplicated routes",
                "Sports/DeFi/CeFi/TradFi via FilterBar dimension, not separate surfaces",
            ],
        ),
    ]
    y = 5.7
    for category, items in outcomes:
        add_text_box(slide, 7.0, y, 1.0, 0.18, category, font_size=9, bold=True, color=ACCENT_GREEN)
        for item in items:
            add_text_box(slide, 8.0, y, 4.5, 0.18, item, font_size=8, color=TEXT_SECONDARY)
            y += 0.25

    # Final thesis
    add_card(slide, 0.5, 7.2, 12.3, 0.25, BG_SECONDARY)
    add_text_box(
        slide,
        0.7,
        7.2,
        11.5,
        0.25,
        "The complexity stays. The navigation becomes legible: hierarchy, lifecycle, cross-linking.",
        font_size=11,
        bold=True,
        color=TEXT_PRIMARY,
        alignment=PP_ALIGN.CENTER,
    )


def slide_09_temporal(prs):
    """Temporal universality — snapshot + time series everywhere."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(
        slide,
        0.5,
        0.4,
        10,
        0.4,
        "Temporal universality — snapshot + time series everywhere",
        font_size=28,
        bold=True,
        color=TEXT_PRIMARY,
    )
    add_text_box(
        slide,
        0.5,
        0.85,
        10,
        0.3,
        "Every view showing state supports two temporal modes. This is an overarching design principle.",
        font_size=11,
        color=TEXT_SECONDARY,
    )

    # Left: Point-in-time snapshot
    add_card(slide, 0.5, 1.4, 6.0, 3.0, BG_ELEVATED)
    add_text_box(
        slide,
        0.7,
        1.5,
        5,
        0.3,
        "Point-in-time snapshot (the time machine)",
        font_size=14,
        bold=True,
        color=ACCENT_AMBER,
    )
    add_text_box(
        slide,
        0.7,
        1.85,
        5.5,
        0.4,
        "Pick any datetime. See the exact state of the world at that moment. Same layout, same drill-down, "
        "just at a different point in time.",
        font_size=9,
        color=TEXT_SECONDARY,
    )

    # AsOfDatetimePicker mockup
    add_card(slide, 0.7, 2.4, 5.5, 0.5, BG_CARD)
    add_text_box(slide, 0.9, 2.45, 1.5, 0.2, "AsOfDatetimePicker", font_size=8, bold=True, color=TEXT_MUTED)
    options = [
        ("Live", ACCENT_GREEN),
        ("1h ago", TEXT_SECONDARY),
        ("Start of day", TEXT_SECONDARY),
        ("Yesterday close", TEXT_SECONDARY),
        ("2026-03-14 14:32", ACCENT_AMBER),
    ]
    x = 0.9
    for label, color in options:
        add_pill(slide, x, 2.65, 1.2 if len(label) > 8 else 0.8, label, color, 7)
        x += 1.05 if len(label) > 8 else 0.85

    add_text_box(
        slide,
        0.7,
        3.05,
        5.5,
        0.2,
        "URL:  ?as_of=2026-03-14T14:32:00Z   (omitted = live, shareable, bookmarkable)",
        font_size=8,
        font_name=FONT_MONO,
        color=TEXT_TERTIARY,
    )

    # Two screenshots side by side: Live vs Historical
    add_card(slide, 0.7, 3.35, 2.6, 0.8, BG_CARD)
    add_pill(slide, 0.8, 3.4, 0.6, "Live", ACCENT_GREEN, 7)
    add_text_box(slide, 0.8, 3.65, 2.3, 0.2, "Firm PnL", font_size=8, color=TEXT_TERTIARY)
    add_text_box(
        slide,
        0.8,
        3.85,
        2.3,
        0.2,
        "+$1.42m",
        font_size=14,
        font_name=FONT_MONO,
        bold=True,
        font_color_override=ACCENT_GREEN,
    )

    add_card(slide, 3.5, 3.35, 2.6, 0.8, BG_CARD)
    add_pill(slide, 3.6, 3.4, 1.3, "2026-03-14 14:32", ACCENT_AMBER, 7)
    add_text_box(slide, 3.6, 3.65, 2.3, 0.2, "Firm PnL", font_size=8, color=TEXT_TERTIARY)
    add_text_box(
        slide,
        3.6,
        3.85,
        2.3,
        0.2,
        "+$0.89m",
        font_size=14,
        font_name=FONT_MONO,
        bold=True,
        font_color_override=ACCENT_AMBER,
    )

    # Right: Time series
    add_card(slide, 6.8, 1.4, 6.0, 3.0, BG_ELEVATED)
    add_text_box(
        slide, 7.0, 1.5, 5, 0.3, "Time series (trajectory embedded inline)", font_size=14, bold=True, color=ACCENT_BLUE
    )
    add_text_box(
        slide,
        7.0,
        1.85,
        5.5,
        0.4,
        "Every KPI card shows a sparkline. Every table metric has a 'show over time' toggle. "
        "The trajectory matters as much as the value.",
        font_size=9,
        color=TEXT_SECONDARY,
    )

    # KPI with sparkline example
    add_card(slide, 7.0, 2.4, 2.5, 1.4, BG_CARD)
    add_text_box(slide, 7.15, 2.5, 2, 0.15, "Firm PnL", font_size=9, color=TEXT_TERTIARY)
    add_text_box(
        slide,
        7.15,
        2.7,
        2,
        0.3,
        "+$1.42m",
        font_size=20,
        font_name=FONT_MONO,
        bold=True,
        font_color_override=ACCENT_GREEN,
    )
    add_text_box(slide, 7.15, 3.0, 2, 0.15, "+0.8% 1d", font_size=8, font_name=FONT_MONO, color=TEXT_MUTED)
    add_sparkline(slide, 7.15, 3.2, 2.1, 0.45, ACCENT_GREEN, "up")
    add_text_box(slide, 7.15, 3.65, 2, 0.12, "1d  1w  1m  3m  1y", font_size=7, color=TEXT_MUTED)

    # TimeSeriesToggle example
    add_card(slide, 9.7, 2.4, 2.8, 1.4, BG_CARD)
    add_text_box(slide, 9.85, 2.5, 2.5, 0.15, "Delta Exposure", font_size=9, color=TEXT_TERTIARY)
    add_text_box(
        slide,
        9.85,
        2.7,
        2.5,
        0.3,
        "$2.4m",
        font_size=20,
        font_name=FONT_MONO,
        bold=True,
        font_color_override=ACCENT_CYAN,
    )
    add_text_box(slide, 9.85, 3.0, 2.5, 0.15, "net long", font_size=8, font_name=FONT_MONO, color=TEXT_MUTED)
    add_sparkline(slide, 9.85, 3.2, 2.3, 0.45, ACCENT_CYAN, "volatile")
    add_text_box(slide, 9.85, 3.65, 2.5, 0.12, "1d  1w  1m  3m  1y", font_size=7, color=TEXT_MUTED)

    # What gets snapshot + time series table
    add_card(slide, 0.5, 4.6, 12.3, 2.8, BG_ELEVATED)
    add_text_box(
        slide, 0.7, 4.7, 10, 0.25, "What gets snapshot + time series", font_size=14, bold=True, color=TEXT_PRIMARY
    )

    domains = [
        (
            "Positions",
            "Exact positions at T",
            "Position size evolution, entry/exit markers",
            "Trading Command Center",
            ACCENT_GREEN,
        ),
        (
            "P&L",
            "Cumulative P&L at T, attribution at T",
            "Equity curve, daily bars, component evolution",
            "Market Intelligence",
            ACCENT_PURPLE,
        ),
        (
            "Risk / Exposure",
            "Delta/funding/basis/greeks at T",
            "Exposure drift, margin trend, LTV evolution",
            "Trading Command Center",
            ACCENT_CYAN,
        ),
        (
            "Alerts",
            "Which alerts were active at T",
            "Alert frequency, resolution time trends",
            "Trading Command Center",
            ACCENT_RED,
        ),
        (
            "Strategy Status",
            "Which strategies live/paused at T",
            "Lifecycle timeline (promoted, paused, killed)",
            "Strategy Analytics",
            ACCENT_BLUE,
        ),
        (
            "Feature Freshness",
            "SLA compliance at T",
            "Freshness degradation, breach history",
            "Trading Command Center",
            ACCENT_AMBER,
        ),
        (
            "Orders / Fills",
            "Orders in flight at T",
            "Order flow over time, fill rate",
            "Market Intelligence",
            ACCENT_PURPLE,
        ),
    ]

    headers = [
        ("Domain", 0.7, 1.5),
        ("Snapshot (as-of T)", 2.2, 2.5),
        ("Time Series (over period)", 4.8, 3.0),
        ("Canonical Home", 8.0, 2.0),
    ]
    for label, x, w in headers:
        add_text_box(slide, x, 5.0, w, 0.2, label, font_size=8, bold=True, color=TEXT_MUTED)

    y = 5.25
    for domain, snapshot, timeseries, home, color in domains:
        add_text_box(slide, 0.7, y, 1.5, 0.18, domain, font_size=9, bold=True, font_color_override=color)
        add_text_box(slide, 2.2, y, 2.5, 0.18, snapshot, font_size=8, color=TEXT_SECONDARY)
        add_text_box(slide, 4.8, y, 3.0, 0.18, timeseries, font_size=8, color=TEXT_SECONDARY)
        add_text_box(slide, 8.0, y, 2.0, 0.18, home, font_size=8, font_color_override=color)
        y += 0.3
